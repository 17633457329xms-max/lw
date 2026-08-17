# -*- coding: utf-8 -*-
"""
第52期 PPT 生成脚本（学术风格，延续第01/09期）
《摘要：论文的"门面"，300字决定第一印象》
干货核心：摘要的功能定位 + 5大要素 + 300字黄金配比 + 正反对比 + 自查清单
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ========== 学术低饱和配色（同第01/09期）==========
BG=RGBColor(0xF6,0xF2,0xEA); PAPER=RGBColor(0xFF,0xFD,0xF7)
DARK=RGBColor(0x2B,0x2B,0x2B); NAVY=RGBColor(0x2C,0x3E,0x50); INK=RGBColor(0x1F,0x3A,0x5F)
RUST=RGBColor(0x8B,0x4A,0x3A); OLIVE=RGBColor(0x5C,0x6B,0x4A); GRAY=RGBColor(0x6E,0x68,0x60)
LINEC=RGBColor(0xC8,0xC0,0xB0); CREAM=RGBColor(0xEF,0xEA,0xDC); WHITE=RGBColor(0xFF,0xFF,0xFF)
SONG='宋体'; HEI='黑体'; KAI='楷体'; SERIF='Times New Roman'

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
blank=prs.slide_layouts[6]

def slide(bg=BG):
    s=prs.slides.add_slide(blank)
    g=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    g.fill.solid(); g.fill.fore_color.rgb=bg; g.line.fill.background(); g.shadow.inherit=False
    return s

def txt(s,l,t,w,h,text,size=18,color=DARK,bold=False,align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,font=HEI,italic=False,lh=None):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    for i,line in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        if lh: p.line_spacing=lh
        r=p.add_run(); r.text=line
        r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=color; r.font.name=font
    return tb

def rect(s,l,t,w,h,fill=None,line=None,line_w=1.0):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(line_w)
    sh.shadow.inherit=False; return sh

def rrect(s,l,t,w,h,fill=None,line=None,line_w=1.0):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(line_w)
    sh.shadow.inherit=False; return sh

def ln(s,x1,y1,x2,y2,color=NAVY,w=1.5):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,x1,y1,x2,y2)
    c.line.color.rgb=color; c.line.width=Pt(w); c.shadow.inherit=False; return c

def oval(s,l,t,w,h,fill=None,line=None,line_w=1.5):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,l,t,w,h)
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(line_w)
    sh.shadow.inherit=False; return sh

def pageno(s,n,total=16):
    txt(s,Inches(11.7),Inches(7.05),Inches(1.4),Inches(0.3),
        f"P. {n:02d}",size=10,color=GRAY,align=PP_ALIGN.RIGHT,font=SERIF)

def header(s,kicker,color=NAVY):
    rect(s,Inches(0.55),Inches(0.45),Inches(0.05),Inches(0.36),fill=color)
    txt(s,Inches(0.72),Inches(0.42),Inches(9),Inches(0.4),kicker,size=12,color=color,font=KAI)
    txt(s,Inches(9.4),Inches(0.42),Inches(3.3),Inches(0.4),
        "论文写作方法论 · 第五十二讲",size=10,color=GRAY,align=PP_ALIGN.RIGHT,font=SONG)

def hline(s,l,t,w,color=LINEC,w_pt=0.75):
    ln(s,l,t,l+w,t,color=color,w=w_pt)

def style_axis(chart,color=GRAY,size=9):
    try:
        chart.category_axis.tick_labels.font.size=Pt(size)
        chart.category_axis.tick_labels.font.color.rgb=color
        chart.category_axis.format.line.color.rgb=LINEC
        if chart.has_value_axis:
            chart.value_axis.tick_labels.font.size=Pt(size)
            chart.value_axis.tick_labels.font.color.rgb=color
            chart.value_axis.format.line.color.rgb=LINEC
            chart.value_axis.major_gridlines.format.line.color.rgb=LINEC
    except Exception: pass

# ========== 简笔画 ==========
def draw_door(s,l,t,w,h,color=NAVY):
    """门面：门框"""
    rect(s,l,t,w,h,fill=PAPER,line=color,line_w=1.5)
    txt(s,l+Emu(int(w*0.08)),t+Emu(int(h*0.1)),Emu(int(w*0.84)),Emu(int(h*0.6)),
        "摘要",size=16,color=INK,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
    oval(s,l+Emu(int(w*0.55)),t+Emu(int(h*0.45)),Emu(int(w*0.18)),Emu(int(w*0.18)),
         line=color,line_w=1.25)
    ln(s,l+Emu(int(w*0.5)),t,l+Emu(int(w*0.5)),t+Emu(int(h*0.9)),color=color,w=1.0)

def draw_eye(s,cx,cy,r,color=NAVY):
    """眼睛：第一印象"""
    oval(s,cx-Emu(int(r*1.1)),cy-Emu(int(r)),Emu(int(r*2.2)),Emu(int(r*2)),line=color,line_w=1.75)
    oval(s,cx-Emu(int(r*0.45)),cy-Emu(int(r*0.45)),Emu(int(r*0.9)),Emu(int(r*0.9)),fill=color,line=color,line_w=1.0)

def draw_pencil(s,l,t,w,color=RUST):
    """铅笔：写作"""
    ln(s,l,t,l+Emu(int(w*0.75)),t,color=color,w=w//6)
    ln(s,l+Emu(int(w*0.75)),t,l+Emu(int(w*0.95)),t-Emu(int(w*0.12)),color=color,w=w//6)
    ln(s,l+Emu(int(w*0.75)),t,l+Emu(int(w*0.95)),t+Emu(int(w*0.12)),color=color,w=w//6)

def draw_watch(s,cx,cy,r,color=NAVY):
    """秒表/时钟：时间与篇幅控制"""
    oval(s,cx-Emu(int(r)),cy-Emu(int(r)),Emu(int(r*2)),Emu(int(r*2)),line=color,line_w=2.0)
    ln(s,cx,cy,cx,cy-Emu(int(r*0.6)),color=color,w=1.5)
    ln(s,cx,cy,cx+Emu(int(r*0.5)),cy+Emu(int(r*0.2)),color=color,w=1.5)

def draw_star(s,cx,cy,r,color=RUST):
    oval(s,cx-Emu(int(r*0.5)),cy-Emu(int(r*0.5)),Emu(int(r)),Emu(int(r)),fill=color,line=color,line_w=1.0)

def draw_check(s,l,t,w,color=OLIVE):
    ln(s,l,t+Emu(int(w*0.5)),l+Emu(int(w*0.3)),t+Emu(int(w*0.9)),color=color,w=w//5)
    ln(s,l+Emu(int(w*0.3)),t+Emu(int(w*0.9)),l+w,t-Emu(int(w*0.1)),color=color,w=w//5)

# ===================================================================
# P01 封面
# ===================================================================
s=slide()
# 左侧：门面（摘要之窗）
draw_door(s,Inches(1.4),Inches(2.2),Inches(2.6),Inches(3.0),color=NAVY)
draw_eye(s,Inches(3.0),Inches(2.0),Inches(0.4),color=RUST)
txt(s,Inches(1.2),Inches(5.6),Inches(3.0),Inches(0.4),
    "300 字 · 决定第一印象",size=12,color=GRAY,align=PP_ALIGN.CENTER,font=KAI)

txt(s,Inches(4.2),Inches(0.9),Inches(8),Inches(0.5),
    "论文写作方法论 · 第五十二讲",size=15,color=RUST,font=KAI)
hline(s,Inches(4.2),Inches(1.45),Inches(3.0),color=RUST,w_pt=1.25)
txt(s,Inches(4.2),Inches(2.0),Inches(8.5),Inches(2.2),
    "摘要：论文的\"门面\"\n300字决定第一印象",size=42,color=INK,bold=True,font=SONG)
txt(s,Inches(4.2),Inches(4.4),Inches(8.5),Inches(0.8),
    "评委最先看到的就是它，你却总在最后才写它",size=21,color=RUST,font=KAI)
hline(s,Inches(0.6),Inches(6.6),Inches(12.1),color=LINEC,w_pt=1.0)
txt(s,Inches(0.6),Inches(6.75),Inches(8),Inches(0.4),
    "摘要的功能定位 —— 第一印象 · 阅读入口 · 检索关键",size=13,color=GRAY,font=KAI)
txt(s,Inches(9.5),Inches(6.75),Inches(3.2),Inches(0.4),
    "Lecture 52 / 100",size=12,color=GRAY,align=PP_ALIGN.RIGHT,font=SERIF)
pageno(s,1)

# ===================================================================
# P02 钩子：摘要之困
# ===================================================================
s=slide()
header(s,"引　言　　摘要之困")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "写摘要时，下面这些情况你中了几条：",size=23,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 左侧简笔：乱稿纸 + 铅笔
draw_pencil(s,Inches(1.0),Inches(2.3),Inches(1.6),color=RUST)
txt(s,Inches(0.8),Inches(3.3),Inches(2.6),Inches(0.4),
    "最后一刻才动笔",size=12,color=RUST,bold=True,align=PP_ALIGN.CENTER,font=KAI)

# 右侧四条困境
plights=[
    ("1","正文写完了才想起摘要，半小时草草凑一段。"),
    ("2","把摘要当成正文缩略版，什么都想说，结果什么都没说清。"),
    ("3","300字里塞了背景、定义、文献回顾，真正的研究内容没多少。"),
    ("4","导师只看摘要就否了全文 —— 门面都没撑住。"),
]
y=Inches(2.0)
for num,c in plights:
    rect(s,Inches(3.9),y+Inches(0.05),Inches(0.04),Inches(0.78),fill=RUST)
    txt(s,Inches(4.05),y+Inches(0.02),Inches(0.45),Inches(0.78),num,size=15,color=RUST,bold=True,anchor=MSO_ANCHOR.MIDDLE,font=SERIF)
    txt(s,Inches(4.6),y,Inches(8.2),Inches(0.82),c,size=14,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE,lh=1.1)
    y+=Inches(1.08)
txt(s,Inches(0.72),Inches(6.55),Inches(12),Inches(0.4),
    "上述困境的共同点：把摘要当\"收尾任务\"，而非\"门面工程\"。",
    size=13,color=RUST,font=KAI)
pageno(s,2)

# ===================================================================
# P03 核心论点：摘要的功能定位
# ===================================================================
s=slide()
header(s,"本　论　　摘要的功能")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "摘要不是论文的\"结尾\"，而是论文的\"入口\"",size=23,color=INK,bold=True,font=SONG)

# 三功能卡片
funcs=[
    ("第一印象","评审、读者最先看到的\n就是它，决定要不要\n继续往下读",NAVY,draw_eye),
    ("阅读入口","读者靠它判断你的\n研究问题、方法、结论\n是否值得精读",RUST,draw_watch),
    ("检索关键","知网/万方索引的就是\n摘要，关键词与检索\n命中全靠它",OLIVE,draw_check),
]
xs=[Inches(0.8),Inches(4.62),Inches(8.44)]
for i,(t,d,col,icon) in enumerate(funcs):
    x=xs[i]
    rrect(s,x,Inches(1.95),Inches(3.6),Inches(3.5),fill=PAPER,line=col,line_w=1.5)
    rect(s,x,Inches(1.95),Inches(3.6),Inches(0.6),fill=col)
    txt(s,x,Inches(1.95),Inches(3.6),Inches(0.6),t,size=16,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
    if icon is draw_eye: icon(s,x+Inches(1.8),Inches(2.95),Inches(0.35),color=col)
    elif icon is draw_watch: icon(s,x+Inches(1.8),Inches(3.0),Inches(0.4),color=col)
    else: draw_check(s,x+Inches(1.45),Inches(2.85),Inches(0.7),color=col)
    txt(s,x+Inches(0.25),Inches(3.6),Inches(3.1),Inches(1.6),d,size=13,color=DARK,font=KAI,lh=1.3,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

rect(s,Inches(0.8),Inches(5.7),Inches(11.73),Inches(1.45),fill=INK)
txt(s,Inches(1.1),Inches(5.8),Inches(11.2),Inches(0.4),
    "核心命题",size=12,color=RGBColor(0xC8,0xC0,0xB0),font=KAI)
txt(s,Inches(1.1),Inches(6.15),Inches(11.2),Inches(0.9),
    "摘要写得不好，等于让评委带着\"不行\"的印象读你的全文。",
    size=19,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER,font=SONG)
pageno(s,3)

# ===================================================================
# P04 数据：评审最先看哪里
# ===================================================================
s=slide()
header(s,"本　论　　评审的阅读路径")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "评审/导师读一篇论文，第一步几乎都是看摘要：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 左侧文字解读
txt(s,Inches(0.8),Inches(2.1),Inches(5.0),Inches(0.4),
    "图1　评审阅读路径分布",size=12,color=INK,bold=True,font=SONG)
txt(s,Inches(0.8),Inches(2.45),Inches(5.0),Inches(0.3),
    "（示意性数据，基于教学观察）",size=10,color=GRAY,font=KAI)
cd=CategoryChartData()
cd.categories=["摘要","目录/框架","引言/绪论","正文抽查","结论"]
cd.add_series("占比",(38,24,18,14,6))
gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.8),Inches(2.9),Inches(5.4),Inches(3.4),cd).chart
gf.has_title=False; gf.has_legend=False
plot=gf.plots[0]; plot.gap_width=80
ser=plot.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb=NAVY
style_axis(gf)

# 右侧解读
txt(s,Inches(6.6),Inches(2.1),Inches(6),Inches(0.4),
    "数据解读",size=14,color=RUST,bold=True,font=SONG)
hline(s,Inches(6.6),Inches(2.55),Inches(5.9))
points=[
    "摘要占近四成 —— 它是被看得最多的部分，权重远超其他。",
    "看摘要的时间往往不到一分钟 —— 必须在极短时间里说清全貌。",
    "摘要不好，后面的内容再好，也容易被先入为主的印象带偏。",
]
y=Inches(2.8)
for p in points:
    rect(s,Inches(6.6),y+Inches(0.05),Inches(0.04),Inches(0.85),fill=RUST)
    txt(s,Inches(6.78),y,Inches(5.7),Inches(0.95),p,size=14,color=DARK,font=KAI,lh=1.15,anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(1.05)
txt(s,Inches(0.72),Inches(6.6),Inches(12),Inches(0.4),
    "结论：摘要投入产出比极高 —— 值得拿出对待全文的认真劲来写。",
    size=14,color=RUST,bold=True,font=SONG)
pageno(s,4)

# ===================================================================
# P05 摘要的五个要素
# ===================================================================
s=slide()
header(s,"本　论　　摘要的五要素")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "一段合格的中文摘要，说清这五件事：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

dims=[
    ("01","问题","研究要解决什么问题？为什么值得解决？一句话点出背景与缺口。",NAVY),
    ("02","方法","用了什么研究设计/数据/工具？让读者信任你的做法。",RUST),
    ("03","结果","得到什么核心发现？用数据说话，不说空话。",RGBColor(0xA8,0x6A,0x3A)),
    ("04","结论","发现说明什么？有什么判断或主张。",OLIVE),
    ("05","意义","对理论/实践有什么价值？（可略写，一两句即可）",GRAY),
]
xs=[Inches(0.8),Inches(3.62),Inches(6.44),Inches(9.26),Inches(11.6)]
for i,(num,t,d,col) in enumerate(dims):
    x=xs[i]
    rrect(s,x,Inches(2.1),Inches(2.65),Inches(3.7),fill=PAPER,line=col,line_w=1.5)
    txt(s,x,Inches(2.3),Inches(2.65),Inches(0.6),num,size=26,color=col,bold=True,align=PP_ALIGN.CENTER,font=SERIF)
    txt(s,x,Inches(3.0),Inches(2.65),Inches(0.5),t,size=18,color=INK,bold=True,align=PP_ALIGN.CENTER,font=SONG)
    hline(s,x+Inches(0.3),Inches(3.55),Inches(2.05),color=col,w_pt=0.75)
    txt(s,x+Inches(0.22),Inches(3.75),Inches(2.2),Inches(1.9),d,size=11,color=DARK,font=KAI,lh=1.2,align=PP_ALIGN.CENTER)
txt(s,Inches(0.72),Inches(6.15),Inches(12),Inches(0.4),
    "核心逻辑：问题 → 方法 → 结果 → 结论 → 意义，一环扣一环，缺一不可。",
    size=13,color=RUST,font=KAI,align=PP_ALIGN.CENTER)
draw_pencil(s,Inches(6.67),Inches(6.6),Inches(1.0),color=RUST)
pageno(s,5)

# ===================================================================
# P06 300字黄金配比（核心方法）
# ===================================================================
s=slide()
header(s,"方　法　　300字黄金配比")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "300字怎么分配？照着这个配比写，不会乱：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 左：配比说明 + 简笔（秒表）
txt(s,Inches(0.8),Inches(2.1),Inches(5.4),Inches(0.4),
    "图2　300字配比示意",size=12,color=INK,bold=True,font=SONG)
cd=CategoryChartData()
cd.categories=["问题(15%)","方法(25%)","结果(35%)","结论(15%)","意义(10%)"]
cd.add_series("占比",(15,25,35,15,10))
pf=s.shapes.add_chart(XL_CHART_TYPE.PIE,Inches(0.8),Inches(2.6),Inches(5.0),Inches(3.6),cd).chart
pf.has_title=False
pf.has_legend=True; pf.legend.position=XL_LEGEND_POSITION.RIGHT; pf.legend.include_in_layout=False
pf.legend.font.size=Pt(10); pf.legend.font.color.rgb=DARK
plot=pf.plots[0]; plot.has_data_labels=True
dl=plot.data_labels; dl.number_format='0%'; dl.number_format_is_linked=False
dl.font.size=Pt(10); dl.font.color.rgb=WHITE; dl.font.bold=True
cols=[NAVY,RUST,RGBColor(0xA8,0x6A,0x3A),OLIVE,GRAY]
for i,pt in enumerate(plot.series[0].points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb=cols[i]
draw_watch(s,Inches(2.9),Inches(6.5),Inches(0.35),color=NAVY)

# 右：逐段写法
txt(s,Inches(6.6),Inches(2.1),Inches(6),Inches(0.4),
    "各段怎么写",size=14,color=RUST,bold=True,font=SONG)
hline(s,Inches(6.6),Inches(2.55),Inches(5.9))
methods=[
    ("问题 · 约45字","一句背景一句缺口：从现状切入，落到\"已有研究尚未解决\"。"),
    ("方法 · 约75字","研究设计一句话：对象/数据/方法/工具，讲清你怎么做。"),
    ("结果 · 约105字","核心发现两三句：最有价值的结论，用数字/方向说清。"),
    ("结论 · 约45字","从结果推出判断：回应开头的问题，给出你的主张。"),
    ("意义 · 约30字","一句价值：对理论或实践有什么用（可选，宁缺毋滥）。"),
]
y=Inches(2.8)
for t,d in methods:
    rect(s,Inches(6.6),y+Inches(0.05),Inches(0.04),Inches(0.72),fill=RUST)
    txt(s,Inches(6.78),y-Inches(0.02),Inches(2.4),Inches(0.4),t,size=12,color=INK,bold=True,font=SONG)
    txt(s,Inches(6.78),y+Inches(0.3),Inches(5.9),Inches(0.42),d,size=11.5,color=DARK,font=KAI,lh=1.05,anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(0.78)
pageno(s,6)

# ===================================================================
# P07 常见误区（正反对照）
# ===================================================================
s=slide()
header(s,"误　区　　六种常见摘要病")
txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),
    "对照自查：你的摘要踩过哪些坑？",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.7),Inches(11.9))

mistakes=[
    ("流水账病","把摘要写成正文缩略版，每章各来一句，没有重心。",RUST),
    ("背景病","前两段都在讲背景和概念，研究本身迟迟不出现。",NAVY),
    ("空泛病","\"取得一定成果\"\"具有积极意义\"——没有具体内容。",RGBColor(0xA8,0x6A,0x3A)),
    ("无结果病","只讲做了什么，不讲做成了什么、得出什么。",OLIVE),
    ("结论病","开头就是结论，读者不知道结论从哪来、对谁有效。",RUST),
    ("超长病","一写就是四五百字，评审根本没耐心看完。",NAVY),
]
xs=[Inches(0.8),Inches(6.83)]
ys=[Inches(2.0),Inches(3.95)]
for i,(t,d,col) in enumerate(mistakes):
    x=xs[i%2]; y=ys[i//2]
    rect(s,x,y,Inches(0.04),Inches(1.7),fill=col)
    txt(s,x+Inches(0.2),y-Inches(0.02),Inches(3.0),Inches(0.4),t,size=15,color=INK,bold=True,font=SONG)
    txt(s,x+Inches(0.2),y+Inches(0.42),Inches(5.4),Inches(1.2),d,size=12,color=DARK,font=KAI,lh=1.15)
txt(s,Inches(0.72),Inches(6.15),Inches(12),Inches(0.4),
    "诊断口诀：摘要只回答五问（问题/方法/结果/结论/意义），其余都是干扰。",
    size=14,color=RUST,bold=True,font=SONG,align=PP_ALIGN.CENTER)
pageno(s,7)

# ===================================================================
# P08 案例：差摘要 vs 好摘要
# ===================================================================
s=slide()
header(s,"案　例　　正反对比")
txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),
    "同一个选题，两种摘要写法，高下立判",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.7),Inches(11.9))

# 左：差摘要
rect(s,Inches(0.8),Inches(2.0),Inches(5.7),Inches(0.55),fill=RUST)
txt(s,Inches(0.8),Inches(2.0),Inches(5.7),Inches(0.55),"差摘要：流水账",
    size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
rrect(s,Inches(0.8),Inches(2.55),Inches(5.7),Inches(3.5),fill=CREAM,line=LINEC,line_w=0.75)
txt(s,Inches(1.05),Inches(2.7),Inches(5.3),Inches(3.2),
    "\"随着经济发展，物流行业日益重要。本文研究了物流企业\n绩效。首先梳理了相关理论，然后分析了现状，提出了对策，\n希望对行业发展有参考意义。\"",
    size=13,color=DARK,font=KAI,lh=1.35,anchor=MSO_ANCHOR.TOP)
txt(s,Inches(1.05),Inches(5.0),Inches(5.3),Inches(0.9),
    "问题：背景占一半、没有方法、没有数据结果、\"参考意义\"落空。",
    size=12,color=RUST,font=KAI,lh=1.15)

# 右：好摘要
rect(s,Inches(6.83),Inches(2.0),Inches(5.7),Inches(0.55),fill=OLIVE)
txt(s,Inches(6.83),Inches(2.0),Inches(5.7),Inches(0.55),"好摘要：五要素齐全",
    size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
rrect(s,Inches(6.83),Inches(2.55),Inches(5.7),Inches(3.5),fill=PAPER,line=OLIVE,line_w=1.0)
txt(s,Inches(7.08),Inches(2.7),Inches(5.3),Inches(3.2),
    "\"针对A股物流企业绩效评估维度单一的问题，本文基于\n2015—2024年32家上市公司面板数据，构建熵权TOPSIS模型\n进行综合评价。研究发现：营运能力权重最高（0.38），数字化\n转型显著提升综合绩效。据此提出聚焦数字化运营的改进\n建议。\"",
    size=13,color=DARK,font=KAI,lh=1.3,anchor=MSO_ANCHOR.TOP)
txt(s,Inches(7.08),Inches(5.3),Inches(5.3),Inches(0.9),
    "亮点：问题→方法→数据→结果→结论→建议，五问全覆盖，全是干货。",
    size=12,color=OLIVE,font=KAI,lh=1.15)

txt(s,Inches(0.72),Inches(6.35),Inches(12),Inches(0.5),
    "对照启示：好摘要让人\"想知道更多\"，差摘要让人\"只想跳过\"。",
    size=14,color=INK,bold=True,font=SONG,align=PP_ALIGN.CENTER)
pageno(s,8)

# ===================================================================
# P09 字数与语言要求
# ===================================================================
s=slide()
header(s,"规　范　　字数与语言")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "摘要的字数与语言，也有硬规矩：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 左：字数规范（柱图）
txt(s,Inches(0.8),Inches(2.1),Inches(5.4),Inches(0.4),
    "图3　各学历摘要字数建议（字，示意）",size=12,color=INK,bold=True,font=SONG)
cd=CategoryChartData()
cd.categories=["本科","硕士","博士"]
cd.add_series("字数",(300,400,500))
cf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.8),Inches(2.6),Inches(5.4),Inches(3.5),cd).chart
cf.has_title=False; cf.has_legend=False
plot=cf.plots[0]; plot.gap_width=60
ser=plot.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb=NAVY
style_axis(cf)
txt(s,Inches(0.8),Inches(6.3),Inches(5.4),Inches(0.6),
    "注：本科约300字，硕士约400字，博士约500字 —— 以学校要求为准。",
    size=11,color=GRAY,font=KAI)

# 右：语言要求
txt(s,Inches(6.6),Inches(2.1),Inches(6),Inches(0.4),
    "语言三要三不要",size=14,color=RUST,bold=True,font=SONG)
hline(s,Inches(6.6),Inches(2.55),Inches(5.9))
lang=[
    ("要客观","用陈述句陈述事实，不用\"我认为\"\"大概\"。"),
    ("要精确","数据、范围、对象交代清楚，不含糊其辞。"),
    ("要独立","单独可读 —— 不看正文也能看懂研究全貌。"),
]
y=Inches(2.85)
for t,d in lang:
    rect(s,Inches(6.6),y+Inches(0.05),Inches(0.04),Inches(1.0),fill=OLIVE)
    txt(s,Inches(6.78),y-Inches(0.02),Inches(2.0),Inches(0.4),t,size=14,color=OLIVE,bold=True,font=SONG)
    txt(s,Inches(6.78),y+Inches(0.35),Inches(5.8),Inches(0.6),d,size=12,color=DARK,font=KAI,lh=1.1)
    y+=Inches(1.15)
txt(s,Inches(0.72),Inches(6.55),Inches(12),Inches(0.4),
    "一句话原则：摘要要能\"独立成文\" —— 抽出来单看也是一篇完整的微论文。",
    size=13,color=RUST,font=KAI,align=PP_ALIGN.CENTER)
pageno(s,9)

# ===================================================================
# P10 关键词选取
# ===================================================================
s=slide()
header(s,"配　件　　关键词怎么选")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "摘要下面的关键词，是检索的命门，别随便凑：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 左：三条原则
txt(s,Inches(0.8),Inches(2.0),Inches(4),Inches(0.4),"三条原则",size=14,color=INK,bold=True,font=SONG)
steps=[
    ("核心概念","你的研究对象/核心变量，务必出现。"),
    ("方法术语","研究方法、模型名称，方便他人检索同类研究。"),
    ("检索友好","用规范学术词，避免口语化、自造词。"),
]
y=Inches(2.5)
for i,(t,d) in enumerate(steps):
    oval(s,Inches(0.85),y+Inches(0.06),Inches(0.4),Inches(0.4),fill=PAPER,line=NAVY,line_w=1.0)
    txt(s,Inches(0.85),y+Inches(0.06),Inches(0.4),Inches(0.4),str(i+1),size=13,color=NAVY,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SERIF)
    txt(s,Inches(1.4),y-Inches(0.02),Inches(1.6),Inches(0.4),t,size=14,color=INK,bold=True,font=SONG)
    txt(s,Inches(3.05),y,Inches(3.6),Inches(0.5),d,size=12,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE,lh=1.1)
    y+=Inches(0.72)

# 右：好关键词示例
txt(s,Inches(7.2),Inches(2.0),Inches(5.4),Inches(0.4),
    "示例（3–5个为宜）",size=13,color=INK,bold=True,font=SONG)
rrect(s,Inches(7.0),Inches(2.5),Inches(5.6),Inches(2.6),fill=PAPER,line=NAVY,line_w=1.0)
txt(s,Inches(7.25),Inches(2.7),Inches(5.1),Inches(2.3),
    "物流企业绩效 · 熵权TOPSIS · 数字化转型 · A股上市\n公司\n\n每个词都应能在摘要中找到对应内容 ——\n关键词是摘要的\"索引\"，不能脱离正文凭空捏造。",
    size=12,color=DARK,font=KAI,lh=1.3,anchor=MSO_ANCHOR.TOP)
txt(s,Inches(7.2),Inches(5.4),Inches(5.4),Inches(0.5),
    "提示：关键词排序按重要性，核心概念放最前。",
    size=11,color=GRAY,font=KAI)
pageno(s,10)

# ===================================================================
# P11 英文摘要
# ===================================================================
s=slide()
header(s,"拓　展　　英文摘要")
txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),
    "英摘别直接逐句硬翻，要\"重写\"而不是\"翻译\"：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.7),Inches(11.9))

# 左：误区 vs 正解
rect(s,Inches(0.8),Inches(2.0),Inches(5.7),Inches(0.55),fill=RUST)
txt(s,Inches(0.8),Inches(2.0),Inches(5.7),Inches(0.55),"误区：中文直译",
    size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
rrect(s,Inches(0.8),Inches(2.55),Inches(5.7),Inches(3.2),fill=CREAM,line=LINEC,line_w=0.75)
txt(s,Inches(1.05),Inches(2.7),Inches(5.3),Inches(2.9),
    "逐句硬译，中英句式纠缠，\n术语不统一，\n读起来像\"翻译腔\"，\n甚至出现中式英语。",
    size=13,color=DARK,font=KAI,lh=1.3,anchor=MSO_ANCHOR.MIDDLE)

rect(s,Inches(6.83),Inches(2.0),Inches(5.7),Inches(0.55),fill=OLIVE)
txt(s,Inches(6.83),Inches(2.0),Inches(5.7),Inches(0.55),"正解：结构重写",
    size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
rrect(s,Inches(6.83),Inches(2.55),Inches(5.7),Inches(3.2),fill=PAPER,line=OLIVE,line_w=1.0)
txt(s,Inches(7.08),Inches(2.7),Inches(5.3),Inches(2.9),
    "按英文学术句式重写，\n用标准术语（如 TOPSIS、panel data），\n时态统一（方法过去时、结果现在时），\n语态简洁（被动/主动搭配）。",
    size=13,color=DARK,font=KAI,lh=1.3,anchor=MSO_ANCHOR.MIDDLE)

txt(s,Inches(0.72),Inches(6.0),Inches(12),Inches(0.5),
    "实用工具：先写中文五要素，再用 DeepL/工具翻底稿，最后人工按学术句式润色。",
    size=14,color=RUST,bold=True,font=SONG,align=PP_ALIGN.CENTER)
pageno(s,11)

# ===================================================================
# P12 写作时机
# ===================================================================
s=slide()
header(s,"时　机　　什么时候写摘要")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "摘要不该\"最后写\"，而应\"边写边改\"：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 时间轴
stages=[
    ("开题前","写一版\"草稿摘要\"","锁定问题与方法，让导师快速判断方向。",NAVY),
    ("正文后","对照正文修订","以实际做的内容为准，删掉没做的事。",RGBColor(0xA8,0x6A,0x3A)),
    ("定稿前","逐字打磨","300字里不含一个废词，数据精确。",RUST),
]
xs=[Inches(0.8),Inches(4.62),Inches(8.44)]
for i,(t,d,note,col) in enumerate(stages):
    x=xs[i]
    oval(s,x+Inches(1.4),Inches(2.15),Inches(0.8),Inches(0.8),fill=PAPER,line=col,line_w=1.5)
    txt(s,x+Inches(1.4),Inches(2.15),Inches(0.8),Inches(0.8),str(i+1),size=24,color=col,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SERIF)
    txt(s,x,Inches(3.15),Inches(3.6),Inches(0.5),t,size=17,color=INK,bold=True,align=PP_ALIGN.CENTER,font=SONG)
    txt(s,x,Inches(3.7),Inches(3.6),Inches(0.5),d,size=14,color=col,bold=True,align=PP_ALIGN.CENTER,font=KAI)
    rrect(s,x+Inches(0.3),Inches(4.35),Inches(3.0),Inches(1.6),fill=PAPER,line=LINEC,line_w=0.75)
    txt(s,x+Inches(0.5),Inches(4.5),Inches(2.6),Inches(1.3),note,size=12,color=DARK,font=KAI,lh=1.25,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if i<2:
        txt(s,x+Inches(3.6),Inches(2.4),Inches(0.4),Inches(0.4),"→",size=16,color=GRAY,align=PP_ALIGN.CENTER,font=SERIF)

txt(s,Inches(0.72),Inches(6.35),Inches(12),Inches(0.5),
    "要诀：摘要写三遍 —— 开题写方向、正文后写事实、定稿前写精品。",
    size=14,color=INK,bold=True,font=SONG,align=PP_ALIGN.CENTER)
pageno(s,12)

# ===================================================================
# P13 自查清单（可截图）
# ===================================================================
s=slide()
header(s,"自　测　　摘要十问",RUST)
txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),
    "写完摘要，用这十个问题逐一核对：",size=23,color=INK,bold=True,font=SONG)
txt(s,Inches(0.72),Inches(1.58),Inches(12),Inches(0.3),
    "〔建议截图保存〕",size=11,color=RUST,font=KAI)
hline(s,Inches(0.72),Inches(2.0),Inches(11.9))
checks=[
    "问题：开头 45 字内是否说清\"要解决什么\"？",
    "方法：是否交代了研究对象、数据与研究方法？",
    "结果：是否有具体数据/方向，而不是\"取得一定成果\"？",
    "结论：是否回应了开头问题，给出明确判断？",
    "意义：是否一句话点出价值（或主动省略）？",
    "字数：是否在 300 字上下（本科）合理范围内？",
    "独立：不看正文，单读摘要能否看懂全貌？",
    "语言：是否客观、精确、无\"我认为\"式主观表述？",
    "关键词：3–5 个，是否覆盖核心概念与方法？",
    "英摘：是否重写而非直译，术语、时态是否统一？",
]
y=Inches(2.3)
for i,c in enumerate(checks):
    col_i=i%2
    x=Inches(0.8) if col_i==0 else Inches(6.83)
    y_=Inches(2.3)+Inches(0.68)*(i//2)
    oval(s,x+Inches(0.1),y_+Inches(0.08),Inches(0.42),Inches(0.42),fill=PAPER,line=INK,line_w=1.0)
    txt(s,x+Inches(0.1),y_+Inches(0.08),Inches(0.42),Inches(0.42),str(i+1),size=13,color=INK,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SERIF)
    rect(s,x+Inches(0.7),y_+Inches(0.08),Inches(0.04),Inches(0.42),fill=RUST)
    txt(s,x+Inches(0.9),y_,Inches(5.0),Inches(0.58),c,size=13,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(0.72),Inches(6.7),Inches(12),Inches(0.4),
    "十问答完，摘要即达标 —— 这就是\"门面\"该有的样子。",
    size=13,color=RUST,bold=True,font=SONG)
pageno(s,13)

# ===================================================================
# P14 衔接：摘要之后
# ===================================================================
s=slide()
header(s,"展　望　　摘要之后")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "摘要写好了，接下来才轮到\"怎么把门面后面的内容撑住\"：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 流程衔接：摘要 → 关键词 → 引言 → 正文
fl=["摘要","关键词","引言/绪论","正文写作"]; fc=[NAVY,NAVY,GRAY,GRAY]
fx=Inches(1.5)
for i in range(4):
    rrect(s,fx,Inches(2.5),Inches(2.2),Inches(0.7),fill=PAPER,line=fc[i],line_w=1.5)
    txt(s,fx,Inches(2.5),Inches(2.2),Inches(0.7),fl[i],size=15,color=fc[i],bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
    if i<3:
        txt(s,fx+Inches(2.2),Inches(2.5),Inches(0.4),Inches(0.7),"→",size=16,color=GRAY,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SERIF)
    fx+=Inches(2.6)
txt(s,Inches(1.5),Inches(3.4),Inches(10),Inches(0.4),
    "本讲解决\"门面\" —— 后续将逐一展开",size=13,color=RUST,font=KAI,align=PP_ALIGN.CENTER)

# 后续预告列表
txt(s,Inches(0.8),Inches(4.2),Inches(12),Inches(0.4),"本系列后续摘要引言篇将讲：",size=14,color=INK,bold=True,font=SONG)
upcoming=[
    "第53讲　摘要的4要素公式，照着写不会错 —— 目的-方法-结果-结论模板",
    "第54讲　关键词怎么选？别随便凑3-5个词 —— 关键词选取规则",
    "第55讲　引言开篇怎么抓人？从背景到问题的过渡 —— 引言漏斗式写法",
    "第58讲　第一章（绪论）整体结构 —— 绪论六件套",
]
y=Inches(4.7)
for u in upcoming:
    rect(s,Inches(0.9),y+Inches(0.05),Inches(0.04),Inches(0.5),fill=NAVY)
    txt(s,Inches(1.1),y,Inches(11.4),Inches(0.55),u,size=13,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(0.6)
pageno(s,14)

# ===================================================================
# P15 结语 + 下期预告
# ===================================================================
s=slide()
header(s,"结　语")
txt(s,Inches(0.72),Inches(1.2),Inches(12),Inches(0.5),"本讲核心命题",size=15,color=RUST,font=KAI)
hline(s,Inches(0.72),Inches(1.8),Inches(3.0),color=RUST,w_pt=1.25)
txt(s,Inches(0.72),Inches(2.1),Inches(12),Inches(1.2),
    "摘要不是收尾，而是门面 ——\n用300字把研究讲清楚，评委才有兴趣进门。",
    size=23,color=INK,bold=True,font=SONG,lh=1.25)
txt(s,Inches(0.72),Inches(4.05),Inches(12),Inches(0.4),"下期预告",size=14,color=RUST,font=KAI)
hline(s,Inches(0.72),Inches(4.5),Inches(11.9))
txt(s,Inches(0.72),Inches(4.75),Inches(12),Inches(0.4),"第五十三讲",size=13,color=GRAY,font=SONG)
txt(s,Inches(0.72),Inches(5.15),Inches(12),Inches(1.0),
    "摘要的4要素公式，照着写不会错\n—— 目的·方法·结果·结论 直接套",
    size=21,color=INK,bold=True,font=SONG)
hline(s,Inches(0.6),Inches(6.85),Inches(12.1))
txt(s,Inches(0.6),Inches(6.95),Inches(12),Inches(0.4),
    "若觉有益，望垂注之 —— 期期相续，将论文一事，讲之通透。",
    size=12,color=GRAY,font=KAI)
pageno(s,15)

out=r"D:\project\code2\temp_pastemd\lw\52\摘要：论文的门面，300字决定第一印象.ppt"
import os
os.makedirs(os.path.dirname(out),exist_ok=True)
prs.save(out)
print("已保存:",out)
print("总页数:",len(prs.slides._sldIdLst))