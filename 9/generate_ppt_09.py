# -*- coding: utf-8 -*-
"""
第09期 PPT 生成脚本（学术风格，延续第01期）
《不会读文献=不会写论文，先学会"挑"文献》
干货核心：文献筛选的优先级判断 —— 相关性/权威性/时效性/类型 四维度 + 筛选漏斗 + 三级分级
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ========== 学术低饱和配色（同第01期）==========
BG=RGBColor(0xF6,0xF2,0xEA); PAPER=RGBColor(0xFF,0xFD,0xF7)
DARK=RGBColor(0x2B,0x2B,0x2B); NAVY=RGBColor(0x2C,0x3E,0x50); INK=RGBColor(0x1F,0x3A,0x5F)
RUST=RGBColor(0x8B,0x4A,0x3A); OLIVE=RGBColor(0x5C,0x6B,0x4A); GRAY=RGBColor(0x6E,0x68,0x60)
LINEC=RGBColor(0xC8,0xC0,0xB0); CREAM=RGBColor(0xEF,0xEA,0xDC); WHITE=RGBColor(0xFF,0xFF,0xFF)
SONG='宋体'; HEI='黑体'; KAI='楷体'; SERIF='Times New Roman'

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
blank=prs.slide_layouts[6]

def slide(bg=BG):
    s=prs.slides.add_slide(blank)
    g=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    g.fill.solid(); g.fill.fore_color.rgb=bg; g.line.fill.background(); g.shadow.inherit=False
    return s

def txt(s,l,t,w,h,text,size=18,color=DARK,bold=False,align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,font=HEI,italic=False,lh=None):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    for i,line in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        if lh: p.line_spacing=lh
        r=p.add_run(); r.text=line
        r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=color; r.font.name=font
    return tb

def rect(s,l,t,w,h,fill=None,line=None,line_w=1.0):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(line_w)
    sh.shadow.inherit=False; return sh

def rrect(s,l,t,w,h,fill=None,line=None,line_w=1.0):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(line_w)
    sh.shadow.inherit=False; return sh

def ln(s,x1,y1,x2,y2,color=NAVY,w=1.5):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,x1,y1,x2,y2)
    c.line.color.rgb=color; c.line.width=Pt(w); c.shadow.inherit=False; return c

def oval(s,l,t,w,h,fill=None,line=None,line_w=1.5):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,l,t,w,h)
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(line_w)
    sh.shadow.inherit=False; return sh

def pageno(s,n,total=16):
    txt(s,Inches(11.7),Inches(7.05),Inches(1.4),Inches(0.3),
        f"P. {n:02d}",size=10,color=GRAY,align=PP_ALIGN.RIGHT,font=SERIF)

def header(s,kicker,color=NAVY):
    rect(s,Inches(0.55),Inches(0.45),Inches(0.05),Inches(0.36),fill=color)
    txt(s,Inches(0.72),Inches(0.42),Inches(9),Inches(0.4),kicker,size=12,color=color,font=KAI)
    txt(s,Inches(9.4),Inches(0.42),Inches(3.3),Inches(0.4),
        "论文写作方法论 · 第九讲",size=10,color=GRAY,align=PP_ALIGN.RIGHT,font=SONG)

def hline(s,l,t,w,color=LINEC,w_pt=0.75):
    ln(s,l,t,l+w,t,color=color,w=w_pt)

def style_axis(chart,color=GRAY,size=9):
    try:
        chart.category_axis.tick_labels.font.size=Pt(size)
        chart.category_axis.tick_labels.font.color.rgb=color
        chart.category_axis.format.line.color.rgb=LINEC
        if chart.has_value_axis:
            chart.value_axis.tick_labels.font.size=Pt(size)
            chart.value_axis.tick_labels.font.color.rgb=color
            chart.value_axis.format.line.color.rgb=LINEC
            chart.value_axis.major_gridlines.format.line.color.rgb=LINEC
    except Exception: pass

# ========== 简笔画 ==========
def draw_doc(s,l,t,w,h,color=NAVY,lines_n=4):
    rect(s,l,t,w,h,fill=PAPER,line=color,line_w=1.0)
    ln(s,l+w-Emu(int(w*0.2)),t,l+w,t+Emu(int(h*0.15)),color=color,w=1.0)
    ln(s,l+w-Emu(int(w*0.2)),t,l+w-Emu(int(w*0.2)),t+Emu(int(h*0.15)),color=color,w=0.75)
    ln(s,l+w-Emu(int(w*0.2)),t+Emu(int(h*0.15)),l+w,t+Emu(int(h*0.15)),color=color,w=0.75)
    for i in range(lines_n):
        ly=t+Emu(int(h*0.32))+Emu(int(h*0.15*i))
        end=0.78 if i%2==0 else 0.6
        ln(s,l+Emu(int(w*0.14)),ly,l+Emu(int(w*end)),ly,color=LINEC,w=0.75)

def draw_stack_docs(s,l,t,n,color=NAVY):
    """一摞文献"""
    for i in range(n):
        draw_doc(s,l+Emu(int(0.08*i*914400)),t+Emu(int(0.1*i*914400)),
                 Inches(1.3),Inches(0.9),color=color,lines_n=3)

def draw_funnel(s,l,t,w,h,color=NAVY):
    """筛选漏斗：倒梯形"""
    ln(s,l,t,l+w,t,color=color,w=1.5)
    ln(s,l,t,l+Emu(int(w*0.36)),t+Emu(int(h*0.7)),color=color,w=1.5)
    ln(s,l+w,t,l+Emu(int(w*0.64)),t+Emu(int(h*0.7)),color=color,w=1.5)
    ln(s,l+Emu(int(w*0.36)),t+Emu(int(h*0.7)),l+Emu(int(w*0.64)),t+Emu(int(h*0.7)),color=color,w=1.5)

def draw_magnifier(s,cx,cy,r,color=NAVY):
    oval(s,cx-Emu(int(r)),cy-Emu(int(r)),Emu(int(r*2)),Emu(int(r*2)),line=color,line_w=2.0)
    ln(s,cx+Emu(int(r*0.7)),cy+Emu(int(r*0.7)),cx+Emu(int(r*1.7)),cy+Emu(int(r*1.7)),color=color,w=3.0)

def draw_star(s,cx,cy,r,color=RUST):
    """五角星标记（核心文献）"""
    oval(s,cx-Emu(int(r*0.5)),cy-Emu(int(r*0.5)),Emu(int(r)),Emu(int(r)),fill=color,line=color,line_w=1.0)

def draw_scales(s,cx,cy,w,color=NAVY):
    """天平：权衡之意"""
    ln(s,cx-Emu(int(w*0.12)),cy+Emu(int(w*0.28)),cx+Emu(int(w*0.12)),cy+Emu(int(w*0.28)),color=color,w=1.5)
    ln(s,cx,cy+Emu(int(w*0.05)),cx-Emu(int(w*0.12)),cy+Emu(int(w*0.28)),color=color,w=1.25)
    ln(s,cx,cy+Emu(int(w*0.05)),cx+Emu(int(w*0.12)),cy+Emu(int(w*0.28)),color=color,w=1.25)
    ln(s,cx,cy-Emu(int(w*0.2)),cx,cy+Emu(int(w*0.05)),color=color,w=1.5)
    ln(s,cx-Emu(int(w*0.5)),cy-Emu(int(w*0.2)),cx+Emu(int(w*0.5)),cy-Emu(int(w*0.2)),color=color,w=1.5)
    for sx in [cx-Emu(int(w*0.5)),cx+Emu(int(w*0.5))]:
        ln(s,sx,cy-Emu(int(w*0.2)),sx,cy-Emu(int(w*0.02)),color=color,w=0.75)
        oval(s,sx-Emu(int(w*0.12)),cy-Emu(int(w*0.02)),Emu(int(w*0.24)),Emu(int(w*0.06)),line=color,line_w=1.0)

# ===================================================================
# P01 封面
# ===================================================================
s=slide()
# 左侧：一摞文献 + 放大镜（挑）
draw_stack_docs(s,Inches(1.0),Inches(2.6),4,color=NAVY)
draw_magnifier(s,Inches(2.6),Inches(2.2),Inches(0.55),color=RUST)
txt(s,Inches(0.8),Inches(4.5),Inches(2.8),Inches(0.4),
    "200 篇 → 精读 20 篇",size=12,color=GRAY,align=PP_ALIGN.CENTER,font=KAI)

txt(s,Inches(4.2),Inches(0.9),Inches(8),Inches(0.5),
    "论文写作方法论 · 第九讲",size=15,color=RUST,font=KAI)
hline(s,Inches(4.2),Inches(1.45),Inches(3.0),color=RUST,w_pt=1.25)
txt(s,Inches(4.2),Inches(2.0),Inches(8.5),Inches(2.2),
    "不会读文献=不会写论文\n先学会\"挑\"文献",size=44,color=INK,bold=True,font=SONG)
txt(s,Inches(4.2),Inches(4.4),Inches(8.5),Inches(0.8),
    "读文献之前，先过一道筛子",size=22,color=RUST,font=KAI)
hline(s,Inches(0.6),Inches(6.6),Inches(12.1),color=LINEC,w_pt=1.0)
txt(s,Inches(0.6),Inches(6.75),Inches(8),Inches(0.4),
    "文献筛选的优先级判断 —— 相关性 · 权威性 · 时效性 · 类型",size=13,color=GRAY,font=KAI)
txt(s,Inches(9.5),Inches(6.75),Inches(3.2),Inches(0.4),
    "Lecture 09 / 100",size=12,color=GRAY,align=PP_ALIGN.RIGHT,font=SERIF)
pageno(s,1)

# ===================================================================
# P02 钩子：痛点场景
# ===================================================================
s=slide()
header(s,"引　言　　文献之困")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "动笔之前，先看看这些文献困境是否常见：",size=23,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 左侧简笔：一摞厚文献 + 发愁人
draw_stack_docs(s,Inches(0.9),Inches(2.0),6,color=NAVY)
txt(s,Inches(0.8),Inches(3.6),Inches(2.6),Inches(0.4),
    "下载 200 篇",size=12,color=RUST,bold=True,align=PP_ALIGN.CENTER,font=KAI)

# 右侧四条困境
plights=[
    ("1","检索结果几百篇，全下回来，真正读完的不到十分之一。"),
    ("2","花一周读完一批文献，发现大半与选题无关，时间白费。"),
    ("3","精读一篇才发现，期刊档次低、观点陈旧，没有参考价值。"),
    ("4","该读的经典没读，不该读的边角料读了一堆，主次颠倒。"),
]
y=Inches(2.0)
for num,c in plights:
    rect(s,Inches(3.9),y+Inches(0.05),Inches(0.04),Inches(0.78),fill=RUST)
    txt(s,Inches(4.05),y+Inches(0.02),Inches(0.45),Inches(0.78),num,size=15,color=RUST,bold=True,anchor=MSO_ANCHOR.MIDDLE,font=SERIF)
    txt(s,Inches(4.6),y,Inches(8.2),Inches(0.82),c,size=14,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE,lh=1.1)
    y+=Inches(1.08)
txt(s,Inches(0.72),Inches(6.55),Inches(12),Inches(0.4),
    "上述困境的共同点：只知\"读\"，不知\"挑\" —— 把精力平均分给了无关文献。",
    size=13,color=RUST,font=KAI)
pageno(s,2)

# ===================================================================
# P03 核心论点：先挑后读
# ===================================================================
s=slide()
header(s,"本　论　　先挑后读")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "一个常被忽视的先后顺序：",size=23,color=INK,bold=True,font=SONG)

# 左：误区
rect(s,Inches(0.8),Inches(1.95),Inches(5.7),Inches(0.55),fill=RUST)
txt(s,Inches(0.8),Inches(1.95),Inches(5.7),Inches(0.55),"认识误区",
    size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
rrect(s,Inches(0.8),Inches(2.5),Inches(5.7),Inches(2.3),fill=CREAM,line=LINEC,line_w=0.75)
txt(s,Inches(1.05),Inches(2.65),Inches(5.3),Inches(2.0),
    "文献 = 读得越多越好\n\n认为下载量越大、阅读量越大，\n论文就越有底气。",
    size=16,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE,lh=1.15)
draw_stack_docs(s,Inches(5.4),Inches(2.8),3,color=RUST)

# 右：正解
rect(s,Inches(6.83),Inches(1.95),Inches(5.7),Inches(0.55),fill=NAVY)
txt(s,Inches(6.83),Inches(1.95),Inches(5.7),Inches(0.55),"本质正解",
    size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
rrect(s,Inches(6.83),Inches(2.5),Inches(5.7),Inches(2.3),fill=PAPER,line=NAVY,line_w=1.0)
txt(s,Inches(7.08),Inches(2.65),Inches(5.3),Inches(2.0),
    "文献 = 先筛后读、精挑细读\n\n读什么比读多少更重要，\n筛选是阅读的前置工序。",
    size=16,color=INK,bold=True,font=KAI,anchor=MSO_ANCHOR.MIDDLE,lh=1.15)
draw_magnifier(s,Inches(11.6),Inches(3.4),Inches(0.5),color=NAVY)

rect(s,Inches(0.8),Inches(5.05),Inches(11.73),Inches(1.6),fill=INK)
txt(s,Inches(1.1),Inches(5.15),Inches(11.2),Inches(0.4),
    "核心命题",size=12,color=RGBColor(0xC8,0xC0,0xB0),font=KAI)
txt(s,Inches(1.1),Inches(5.55),Inches(11.2),Inches(1.0),
    "不会挑文献，等于不会读文献 —— 筛选是阅读的前提，而非可有可无的步骤。",
    size=19,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER,font=SONG)
pageno(s,3)

# ===================================================================
# P04 文献检索痛点分布柱图
# ===================================================================
s=slide()
header(s,"本　论　　文献检索之痛")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "学生在文献环节最常卡在哪里：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 左侧文字解读
txt(s,Inches(0.8),Inches(2.1),Inches(5.0),Inches(0.4),
    "图1　文献环节常见痛点分布",size=12,color=INK,bold=True,font=SONG)
txt(s,Inches(0.8),Inches(2.45),Inches(5.0),Inches(0.3),
    "（示意性数据，基于教学观察）",size=10,color=GRAY,font=KAI)
cd=CategoryChartData()
cd.categories=["不知挑哪些","读完即忘","读不懂","找不到文献","重复阅读"]
cd.add_series("占比",(35,26,20,12,7))
gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.8),Inches(2.9),Inches(5.4),Inches(3.4),cd).chart
gf.has_title=False; gf.has_legend=False
plot=gf.plots[0]; plot.gap_width=80
ser=plot.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb=NAVY
style_axis(gf)

# 右侧解读
txt(s,Inches(6.6),Inches(2.1),Inches(6),Inches(0.4),
    "数据解读",size=14,color=RUST,bold=True,font=SONG)
hline(s,Inches(6.6),Inches(2.55),Inches(5.9))
points=[
    "首位\"不知挑哪些\"占 35% —— 即筛选能力缺失，是本讲要解决的核心。",
    "\"读完即忘\"占 26% —— 表面是记忆问题，实则是没筛、乱读导致信息过载。",
    "二者合计逾六成，根因皆指向\"缺乏筛选\"。",
]
y=Inches(2.8)
for p in points:
    rect(s,Inches(6.6),y+Inches(0.05),Inches(0.04),Inches(0.85),fill=RUST)
    txt(s,Inches(6.78),y,Inches(5.7),Inches(0.95),p,size=14,color=DARK,font=KAI,lh=1.15,anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(1.05)
txt(s,Inches(0.72),Inches(6.6),Inches(12),Inches(0.4),
    "结论：与其追求阅读量，不如先建立一套筛选标准。",
    size=14,color=RUST,bold=True,font=SONG)
pageno(s,4)

# ===================================================================
# P05 挑文献的四个维度
# ===================================================================
s=slide()
header(s,"本　论　　筛选的四维度")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "判断一篇文献值不值得读，看这四个维度：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

dims=[
    ("01","相关性","与你的研究问题是否直接相关。这是第一道闸门，不相关直接淘汰。",NAVY),
    ("02","权威性","来源是否可靠：期刊层次、被引量、作者学术地位。",RUST),
    ("03","时效性","发表时间是否契合：是追最新进展，还是夯实经典基础。",RGBColor(0xA8,0x6A,0x3A)),
    ("04","类型","文献性质：综述/实证/理论，不同类型承担不同功能。",OLIVE),
]
xs=[Inches(0.8),Inches(3.62),Inches(6.44),Inches(9.26)]
for i,(num,t,d,col) in enumerate(dims):
    x=xs[i]
    rrect(s,x,Inches(2.1),Inches(2.65),Inches(3.5),fill=PAPER,line=col,line_w=1.5)
    txt(s,x,Inches(2.3),Inches(2.65),Inches(0.6),num,size=28,color=col,bold=True,align=PP_ALIGN.CENTER,font=SERIF)
    txt(s,x,Inches(3.0),Inches(2.65),Inches(0.5),t,size=18,color=INK,bold=True,align=PP_ALIGN.CENTER,font=SONG)
    hline(s,x+Inches(0.3),Inches(3.55),Inches(2.05),color=col,w_pt=0.75)
    txt(s,x+Inches(0.22),Inches(3.75),Inches(2.2),Inches(1.7),d,size=12,color=DARK,font=KAI,lh=1.2,align=PP_ALIGN.CENTER)
txt(s,Inches(0.72),Inches(6.0),Inches(12),Inches(0.4),
    "四维度并非并列平均 —— 相关性是前提，不满足则其余三项免谈。",
    size=13,color=RUST,font=KAI,align=PP_ALIGN.CENTER)
# 底部简笔：天平（权衡）
draw_scales(s,Inches(6.67),Inches(6.75),Inches(1.2),color=NAVY)
pageno(s,5)

# ===================================================================
# P06 维度一：相关性
# ===================================================================
s=slide()
header(s,"维度一　　相关性")
txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),
    "相关性 —— 第一道闸门，淘汰 70% 的无关文献",size=22,color=NAVY,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.7),Inches(11.9))

# 左：三步快速判断
txt(s,Inches(0.8),Inches(2.0),Inches(4),Inches(0.4),"三步快速判断",size=14,color=INK,bold=True,font=SONG)
steps=[
    ("看标题","标题关键词与你的研究问题是否匹配。"),
    ("看摘要","摘要的研究问题、方法、结论是否与你的相关。"),
    ("看结论","结论是否能为你的论证提供支撑或反例。"),
]
y=Inches(2.5)
for i,(t,d) in enumerate(steps):
    oval(s,Inches(0.85),y+Inches(0.06),Inches(0.4),Inches(0.4),fill=PAPER,line=NAVY,line_w=1.0)
    txt(s,Inches(0.85),y+Inches(0.06),Inches(0.4),Inches(0.4),str(i+1),size=13,color=NAVY,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SERIF)
    txt(s,Inches(1.4),y-Inches(0.02),Inches(1.4),Inches(0.4),t,size=14,color=INK,bold=True,font=SONG)
    txt(s,Inches(2.85),y,Inches(4.0),Inches(0.5),d,size=12,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE,lh=1.1)
    y+=Inches(0.72)

# 右：相关性等级判断 + 饼图
txt(s,Inches(7.2),Inches(2.0),Inches(5.4),Inches(0.4),
    "图2　按相关性分级的文献占比（示意）",size=12,color=INK,bold=True,font=SONG)
cd=CategoryChartData()
cd.categories=["高相关(精读)","中相关(选读)","低相关(备查)","无关(淘汰)"]
cd.add_series("占比",(15,25,20,40))
pf=s.shapes.add_chart(XL_CHART_TYPE.PIE,Inches(7.0),Inches(2.5),Inches(5.6),Inches(3.5),cd).chart
pf.has_title=False
pf.has_legend=True; pf.legend.position=XL_LEGEND_POSITION.RIGHT; pf.legend.include_in_layout=False
pf.legend.font.size=Pt(11); pf.legend.font.color.rgb=DARK
plot=pf.plots[0]; plot.has_data_labels=True
dl=plot.data_labels; dl.number_format='0%'; dl.number_format_is_linked=False
dl.font.size=Pt(11); dl.font.color.rgb=WHITE; dl.font.bold=True
cols=[NAVY,RGBColor(0xA8,0x6A,0x3A),OLIVE,GRAY]
for i,pt in enumerate(plot.series[0].points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb=cols[i]
txt(s,Inches(7.2),Inches(6.1),Inches(5.4),Inches(0.5),
    "约四成文献经此闸门即被淘汰 —— 读完标题摘要即可判断。",
    size=12,color=RUST,font=KAI)
pageno(s,6)

# ===================================================================
# P07 维度二：权威性
# ===================================================================
s=slide()
header(s,"维度二　　权威性")
txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),
    "权威性 —— 判断来源是否可靠的三条线索",size=22,color=RUST,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.7),Inches(11.9))

# 左：三条线索
clues=[
    ("期刊层次","CSSCI/北大核心优先，普通期刊谨慎，非学术来源慎用。"),
    ("被引量","高被引说明学界认可，但需警惕\"高引低质\"的自引泡沫。"),
    ("作者地位","该领域活跃学者、有系列成果的作者，文献更值得信赖。"),
]
y=Inches(2.1)
for t,d in clues:
    rect(s,Inches(0.8),y,Inches(0.05),Inches(1.25),fill=RUST)
    txt(s,Inches(1.0),y-Inches(0.02),Inches(2.0),Inches(0.5),t,size=15,color=INK,bold=True,font=SONG)
    txt(s,Inches(1.0),y+Inches(0.45),Inches(5.6),Inches(0.8),d,size=13,color=DARK,font=KAI,lh=1.15)
    y+=Inches(1.45)

# 右：期刊层次参考价值柱图
txt(s,Inches(7.0),Inches(2.1),Inches(5.6),Inches(0.4),
    "图3　不同期刊层次的参考价值（示意）",size=12,color=INK,bold=True,font=SONG)
cd=CategoryChartData()
cd.categories=["顶刊/权威","CSSCI/核心","普通期刊","非学术来源"]
cd.add_series("价值",(95,80,45,15))
cf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(7.0),Inches(2.6),Inches(5.6),Inches(3.3),cd).chart
cf.has_title=False; cf.has_legend=False
plot=cf.plots[0]; plot.gap_width=60
ser=plot.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb=RUST
style_axis(cf)
txt(s,Inches(7.0),Inches(6.0),Inches(5.6),Inches(0.6),
    "提示：权威性是\"加权\"而非\"一票否决\" ——\n普通期刊亦可作参考，但不宜作主要依据。",
    size=12,color=GRAY,font=KAI,lh=1.15)
pageno(s,7)

# ===================================================================
# P08 维度三：时效性
# ===================================================================
s=slide()
header(s,"维度三　　时效性")
txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),
    "时效性 —— 经典与前沿，如何配比",size=22,color=RGBColor(0xA8,0x6A,0x3A),bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.7),Inches(11.9))

# 左：两类文献
txt(s,Inches(0.8),Inches(2.0),Inches(5.8),Inches(0.4),"两类文献各有所长",size=14,color=INK,bold=True,font=SONG)
rrect(s,Inches(0.8),Inches(2.5),Inches(2.8),Inches(2.4),fill=CREAM,line=NAVY,line_w=1.0)
txt(s,Inches(0.95),Inches(2.65),Inches(2.5),Inches(0.4),"经典文献",size=15,color=NAVY,bold=True,align=PP_ALIGN.CENTER,font=SONG)
txt(s,Inches(0.95),Inches(3.1),Inches(2.5),Inches(1.7),
    "奠基性、高被引\n\n奠定理论基础\n提供学科脉络",size=12,color=DARK,font=KAI,lh=1.2,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
rrect(s,Inches(3.8),Inches(2.5),Inches(2.8),Inches(2.4),fill=PAPER,line=RGBColor(0xA8,0x6A,0x3A),line_w=1.0)
txt(s,Inches(3.95),Inches(2.65),Inches(2.5),Inches(0.4),"最新文献",size=15,color=RGBColor(0xA8,0x6A,0x3A),bold=True,align=PP_ALIGN.CENTER,font=SONG)
txt(s,Inches(3.95),Inches(3.1),Inches(2.5),Inches(1.7),
    "近3–5年\n\n反映研究前沿\n填补研究空白",size=12,color=DARK,font=KAI,lh=1.2,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# 右：时间结构配比饼图
txt(s,Inches(7.0),Inches(2.0),Inches(5.6),Inches(0.4),
    "图4　建议的文献时间结构（示意）",size=12,color=INK,bold=True,font=SONG)
cd=CategoryChartData()
cd.categories=["近5年","5–10年","10年以上经典"]
cd.add_series("占比",(50,30,20))
pf=s.shapes.add_chart(XL_CHART_TYPE.PIE,Inches(7.0),Inches(2.5),Inches(5.6),Inches(3.4),cd).chart
pf.has_title=False
pf.has_legend=True; pf.legend.position=XL_LEGEND_POSITION.RIGHT; pf.legend.include_in_layout=False
pf.legend.font.size=Pt(11); pf.legend.font.color.rgb=DARK
plot=pf.plots[0]; plot.has_data_labels=True
dl=plot.data_labels; dl.number_format='0%'; dl.number_format_is_linked=False
dl.font.size=Pt(11); dl.font.color.rgb=WHITE; dl.font.bold=True
cols=[RGBColor(0xA8,0x6A,0x3A),NAVY,GRAY]
for i,pt in enumerate(plot.series[0].points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb=cols[i]
txt(s,Inches(7.0),Inches(6.0),Inches(5.6),Inches(0.5),
    "原则：以新为主、以旧为基 —— 新文献占半数，经典文献打底。",
    size=12,color=RGBColor(0xA8,0x6A,0x3A),font=KAI)
pageno(s,8)

# ===================================================================
# P09 维度四：类型
# ===================================================================
s=slide()
header(s,"维度四　　文献类型")
txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),
    "类型 —— 综述、实证、理论，各司其职",size=22,color=OLIVE,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.7),Inches(11.9))

types=[
    ("综述类","快速建立领域地图，了解研究脉络与争议。","先读，定方向",NAVY),
    ("实证类","提供方法范例与数据参照，支撑你的研究设计。","中读，学方法",RGBColor(0xA8,0x6A,0x3A)),
    ("理论类","提供分析框架与概念工具，提升论证深度。","后读，强论证",OLIVE),
]
xs=[Inches(0.8),Inches(4.62),Inches(8.44)]
for i,(t,d,role,col) in enumerate(types):
    x=xs[i]
    rrect(s,x,Inches(2.1),Inches(3.6),Inches(3.6),fill=PAPER,line=col,line_w=1.5)
    rect(s,x,Inches(2.1),Inches(3.6),Inches(0.6),fill=col)
    txt(s,x,Inches(2.1),Inches(3.6),Inches(0.6),t,size=16,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
    txt(s,x+Inches(0.25),Inches(2.95),Inches(3.1),Inches(1.8),d,size=13,color=DARK,font=KAI,lh=1.2,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    hline(s,x+Inches(0.6),Inches(4.85),Inches(2.4),color=col,w_pt=0.75)
    txt(s,x,Inches(5.0),Inches(3.6),Inches(0.5),role,size=13,color=col,bold=True,align=PP_ALIGN.CENTER,font=SONG)
txt(s,Inches(0.72),Inches(6.2),Inches(12),Inches(0.5),
    "建议配比：综述打底、实证为体、理论为骨 —— 三类搭配，方成完整文献体系。",
    size=14,color=OLIVE,bold=True,font=SONG,align=PP_ALIGN.CENTER)
pageno(s,9)

# ===================================================================
# P10 筛选漏斗（核心方法）
# ===================================================================
s=slide()
header(s,"方　法　　筛选漏斗")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "把筛选做成一个漏斗 —— 四级递减，层层过滤",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 左：漏斗简笔 + 层级标注
fx=Inches(1.2); fy=Inches(2.3); fw=Inches(3.6); fh=Inches(3.6)
draw_funnel(s,fx,fy,fw,fh,color=NAVY)
# 四层标注
levels=[
    ("检索结果","200+ 篇",NAVY),
    ("粗筛（标题摘要）","约 60 篇",RGBColor(0xA8,0x6A,0x3A)),
    ("精筛（全文）","约 30 篇",RUST),
    ("精读","约 20 篇",OLIVE),
]
ly=fy
for i,(name,num,col) in enumerate(levels):
    txt(s,Inches(5.0),ly-Inches(0.05),Inches(2.5),Inches(0.4),name,size=14,color=col,bold=True,font=SONG)
    txt(s,Inches(7.6),ly-Inches(0.05),Inches(2.0),Inches(0.4),num,size=14,color=DARK,font=KAI)
    if i<3:
        ln(s,Inches(4.85),ly+Inches(0.15),Inches(9.6),ly+Inches(0.15),color=LINEC,w=0.5)
    ly+=Inches(0.95)

# 右：漏斗各层做法
txt(s,Inches(9.8),Inches(2.0),Inches(3.2),Inches(0.4),"各层做法",size=13,color=INK,bold=True,font=SONG)
methods=[
    "检索：用主题词+布尔逻辑扩检",
    "粗筛：读标题摘要，相关性筛选",
    "精筛：通读全文，权威性时效性筛选",
    "精读：逐篇研读，做笔记建联系",
]
y=Inches(2.45)
for m in methods:
    rrect(s,Inches(9.8),y,Inches(3.2),Inches(0.8),fill=CREAM,line=LINEC,line_w=0.75)
    txt(s,Inches(9.95),y,Inches(2.9),Inches(0.8),m,size=12,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE,lh=1.1)
    y+=Inches(0.9)

txt(s,Inches(0.72),Inches(6.55),Inches(12),Inches(0.4),
    "要旨：每层只筛掉，不回头读 —— 越往后越少越精，避免\"全都要\"的陷阱。",
    size=13,color=RUST,font=KAI,align=PP_ALIGN.CENTER)
pageno(s,10)

# ===================================================================
# P11 三级文献分级法
# ===================================================================
s=slide()
header(s,"方　法　　三级分级法")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "筛完后，给文献分三个等级，区别对待",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

grades=[
    ("A 级","必读","核心文献，反复研读、做详注、建引用。",NAVY,"约 15%"),
    ("B 级","选读","相关文献，通读要点、记观点、备引用。",RGBColor(0xA8,0x6A,0x3A),"约 25%"),
    ("C 级","备查","边缘文献，知其存在即可，用时再查。",GRAY,"约 60%"),
]
xs=[Inches(0.8),Inches(4.62),Inches(8.44)]
for i,(g,role,d,col,pct) in enumerate(grades):
    x=xs[i]
    rrect(s,x,Inches(2.1),Inches(3.6),Inches(3.8),fill=PAPER,line=col,line_w=1.5)
    rect(s,x,Inches(2.1),Inches(3.6),Inches(0.7),fill=col)
    txt(s,x,Inches(2.1),Inches(3.6),Inches(0.7),g,size=18,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
    txt(s,x,Inches(2.95),Inches(3.6),Inches(0.5),role,size=16,color=col,bold=True,align=PP_ALIGN.CENTER,font=SONG)
    txt(s,x+Inches(0.25),Inches(3.5),Inches(3.1),Inches(1.6),d,size=13,color=DARK,font=KAI,lh=1.2,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    hline(s,x+Inches(0.6),Inches(5.15),Inches(2.4),color=col,w_pt=0.75)
    txt(s,x,Inches(5.3),Inches(3.6),Inches(0.5),pct,size=14,color=col,bold=True,align=PP_ALIGN.CENTER,font=SERIF)
txt(s,Inches(0.72),Inches(6.3),Inches(12),Inches(0.5),
    "精力分配应倒置：A 级文献虽少，却要占六成以上阅读时间。",
    size=14,color=RUST,bold=True,font=SONG,align=PP_ALIGN.CENTER)
pageno(s,11)

# ===================================================================
# P12 案例：正反对比
# ===================================================================
s=slide()
header(s,"案　例　　正反对比")
txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),
    "两种文献处理方式，结果天壤之别",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.7),Inches(11.9))

# 左：反面
rect(s,Inches(0.8),Inches(2.0),Inches(5.7),Inches(0.55),fill=RUST)
txt(s,Inches(0.8),Inches(2.0),Inches(5.7),Inches(0.55),"反面：来者不拒",
    size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
rrect(s,Inches(0.8),Inches(2.55),Inches(5.7),Inches(3.4),fill=CREAM,line=LINEC,line_w=0.75)
txt(s,Inches(1.05),Inches(2.7),Inches(5.3),Inches(3.1),
    "做法：检索结果全下，逐篇通读。\n\n结果：耗时两周，读完一片茫然；\n动笔时发现，有用的不到五篇，\n且彼此观点冲突，无法组织。\n\n症结：无筛选，无分级，无重点。",
    size=13,color=DARK,font=KAI,lh=1.25,anchor=MSO_ANCHOR.TOP)
draw_stack_docs(s,Inches(5.4),Inches(4.6),3,color=RUST)

# 右：正面
rect(s,Inches(6.83),Inches(2.0),Inches(5.7),Inches(0.55),fill=OLIVE)
txt(s,Inches(6.83),Inches(2.0),Inches(5.7),Inches(0.55),"正面：先筛后读",
    size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
rrect(s,Inches(6.83),Inches(2.55),Inches(5.7),Inches(3.4),fill=PAPER,line=OLIVE,line_w=1.0)
txt(s,Inches(7.08),Inches(2.7),Inches(5.3),Inches(3.1),
    "做法：四维度筛 → 漏斗过滤 → ABC 分级。\n\n结果：三天筛完，锁定 20 篇精读；\n动笔时脉络清晰，A 级文献\n直接支撑各章论证。\n\n要旨：读得少，但读得准、读得透。",
    size=13,color=DARK,font=KAI,lh=1.25,anchor=MSO_ANCHOR.TOP)
draw_magnifier(s,Inches(11.6),Inches(4.9),Inches(0.45),color=OLIVE)

txt(s,Inches(0.72),Inches(6.25),Inches(12),Inches(0.5),
    "对比启示：文献功夫在\"筛\"不在\"读\" —— 少而精，胜过多而乱。",
    size=14,color=INK,bold=True,font=SONG,align=PP_ALIGN.CENTER)
pageno(s,12)

# ===================================================================
# P13 精读文献数量建议（柱图）
# ===================================================================
s=slide()
header(s,"参　照　　精读数量建议")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "不同学历层次的精读文献量参考：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 左：柱图
txt(s,Inches(0.8),Inches(2.1),Inches(5.4),Inches(0.4),
    "图5　精读文献数量建议（篇，示意）",size=12,color=INK,bold=True,font=SONG)
cd=CategoryChartData()
cd.categories=["课程论文","本科毕业","硕士毕业","博士毕业"]
cd.add_series("篇数",(8,20,40,80))
cf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.8),Inches(2.6),Inches(5.4),Inches(3.5),cd).chart
cf.has_title=False; cf.has_legend=False
plot=cf.plots[0]; plot.gap_width=60
ser=plot.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb=NAVY
style_axis(cf)

# 右：说明
txt(s,Inches(6.6),Inches(2.1),Inches(6),Inches(0.4),"数量说明",size=14,color=RUST,bold=True,font=SONG)
hline(s,Inches(6.6),Inches(2.55),Inches(5.9))
notes=[
    "此为\"精读\"数量，不含粗筛与备查文献。",
    "精读量随学历递增，博士需建立系统文献体系。",
    "数量是参考而非硬指标 —— 以\"读懂、能用\"为准。",
    "宁少而透，勿多而浮：读懂 20 篇胜过翻过 200 篇。",
]
y=Inches(2.8)
for n in notes:
    rect(s,Inches(6.6),y+Inches(0.05),Inches(0.04),Inches(0.8),fill=RUST)
    txt(s,Inches(6.78),y,Inches(5.7),Inches(0.85),n,size=13,color=DARK,font=KAI,lh=1.15,anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(0.92)
pageno(s,13)

# ===================================================================
# P14 自查清单
# ===================================================================
s=slide()
header(s,"自　测　　筛选五问",RUST)
txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),
    "筛文献时，用这五个问题逐一判断：",size=23,color=INK,bold=True,font=SONG)
txt(s,Inches(0.72),Inches(1.58),Inches(12),Inches(0.3),
    "〔建议截图保存〕",size=11,color=RUST,font=KAI)
hline(s,Inches(0.72),Inches(2.0),Inches(11.9))
checks=[
    "相关性：这篇文献的研究问题，与我的问题是否直接相关？",
    "权威性：来源期刊层次如何？被引量是否足以佐证其价值？",
    "时效性：发表时间是否契合我的需要（前沿 or 经典）？",
    "类型：它是综述、实证还是理论？对我的论文起什么作用？",
    "定位：它该进 A 级（必读）、B 级（选读）还是 C 级（备查）？",
]
y=Inches(2.35)
for i,c in enumerate(checks):
    oval(s,Inches(0.9),y+Inches(0.08),Inches(0.42),Inches(0.42),fill=PAPER,line=INK,line_w=1.0)
    txt(s,Inches(0.9),y+Inches(0.08),Inches(0.42),Inches(0.42),str(i+1),size=14,color=INK,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SERIF)
    rect(s,Inches(1.5),y+Inches(0.08),Inches(0.04),Inches(0.42),fill=RUST)
    txt(s,Inches(1.7),y,Inches(10.8),Inches(0.58),c,size=15,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(0.7)
txt(s,Inches(0.72),Inches(6.7),Inches(12),Inches(0.4),
    "五问答完，每篇文献即有明确去留 —— 这才是\"会挑\"。",
    size=13,color=RUST,bold=True,font=SONG)
pageno(s,14)

# ===================================================================
# P15 衔接：挑完之后怎么读
# ===================================================================
s=slide()
header(s,"展　望　　挑完之后")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "筛完文献，接下来才是\"怎么读\"：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 流程衔接：挑 → 读 → 综 → 用
fl=["挑文献","读文献","写综述","用于论证"]; fc=[NAVY,NAVY,GRAY,GRAY]
fx=Inches(1.5)
for i in range(4):
    rrect(s,fx,Inches(2.5),Inches(2.2),Inches(0.7),fill=PAPER,line=fc[i],line_w=1.5)
    txt(s,fx,Inches(2.5),Inches(2.2),Inches(0.7),fl[i],size=15,color=fc[i],bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
    if i<3:
        txt(s,fx+Inches(2.2),Inches(2.5),Inches(0.4),Inches(0.7),"→",size=16,color=GRAY,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SERIF)
    fx+=Inches(2.6)
txt(s,Inches(1.5),Inches(3.4),Inches(10),Inches(0.4),
    "本讲解决第①步\"挑\" —— 后续将逐一展开",size=13,color=RUST,font=KAI,align=PP_ALIGN.CENTER)

# 后续预告列表
txt(s,Inches(0.8),Inches(4.2),Inches(12),Inches(0.4),"本系列后续文献篇将讲：",size=14,color=INK,bold=True,font=SONG)
upcoming=[
    "第10讲　知网高效检索的7个技巧 —— 检索环节的扩检与精检",
    "第12讲　一天读10篇文献的方法：三遍阅读法 —— 精读环节的分层阅读",
    "第16讲　读文献要不要做笔记？怎么做才不白读 —— 阅读的沉淀",
    "第24讲　文献综述不是文献堆砌 —— 从读到综的转化",
]
y=Inches(4.7)
for u in upcoming:
    rect(s,Inches(0.9),y+Inches(0.05),Inches(0.04),Inches(0.5),fill=NAVY)
    txt(s,Inches(1.1),y,Inches(11.4),Inches(0.55),u,size=13,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(0.6)
pageno(s,15)

# ===================================================================
# P16 结语 + 下期预告
# ===================================================================
s=slide()
header(s,"结　语")
txt(s,Inches(0.72),Inches(1.2),Inches(12),Inches(0.5),"本讲核心命题",size=15,color=RUST,font=KAI)
hline(s,Inches(0.72),Inches(1.8),Inches(3.0),color=RUST,w_pt=1.25)
txt(s,Inches(0.72),Inches(2.1),Inches(12),Inches(1.2),
    "不会挑文献，等于不会读文献 ——\n先过筛子，再开卷。",
    size=24,color=INK,bold=True,font=SONG,lh=1.2)
txt(s,Inches(0.72),Inches(4.05),Inches(12),Inches(0.4),"下期预告",size=14,color=RUST,font=KAI)
hline(s,Inches(0.72),Inches(4.5),Inches(11.9))
txt(s,Inches(0.72),Inches(4.75),Inches(12),Inches(0.4),"第　十　讲",size=13,color=GRAY,font=SONG)
txt(s,Inches(0.72),Inches(5.15),Inches(12),Inches(1.0),
    "知网高效检索的7个技巧\n—— 别再只会输关键词",
    size=21,color=INK,bold=True,font=SONG)
hline(s,Inches(0.6),Inches(6.85),Inches(12.1))
txt(s,Inches(0.6),Inches(6.95),Inches(12),Inches(0.4),
    "若觉有益，望垂注之 —— 期期相续，将论文一事，讲之通透。",
    size=12,color=GRAY,font=KAI)
pageno(s,16)

out=r"D:\project\code2\temp_pastemd\lw\9\不会读文献先学会挑文献.ppt"
prs.save(out)
print("已保存:",out)
print("总页数:",len(prs.slides._sldIdLst))