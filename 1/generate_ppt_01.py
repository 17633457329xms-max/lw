# -*- coding: utf-8 -*-
"""
第01期 PPT 生成脚本（学术化重制 v2）
《毕业论文到底在考你什么？90%的人从第一步就错了》
要求：现代学术规范用语（非文言）、低饱和配色、真实可编辑数据图表、精确分区无重叠
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ========== 学术低饱和配色 ==========
BG     = RGBColor(0xF6, 0xF2, 0xEA)
PAPER  = RGBColor(0xFF, 0xFD, 0xF7)
DARK   = RGBColor(0x2B, 0x2B, 0x2B)
NAVY   = RGBColor(0x2C, 0x3E, 0x50)
INK    = RGBColor(0x1F, 0x3A, 0x5F)
RUST   = RGBColor(0x8B, 0x4A, 0x3A)
OLIVE  = RGBColor(0x5C, 0x6B, 0x4A)
GRAY   = RGBColor(0x6E, 0x68, 0x60)
LINEC  = RGBColor(0xC8, 0xC0, 0xB0)
CREAM  = RGBColor(0xEF, 0xEA, 0xDC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PALEB  = RGBColor(0xDC, 0xE3, 0xEB)   # 浅蓝数据色
PALEO  = RGBColor(0xDD, 0xE2, 0xD2)   # 浅橄榄数据色

SONG='宋体'; HEI='黑体'; KAI='楷体'; SERIF='Times New Roman'

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]

def slide(bg=BG):
    s = prs.slides.add_slide(blank)
    g = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    g.fill.solid(); g.fill.fore_color.rgb=bg; g.line.fill.background(); g.shadow.inherit=False
    return s

def txt(s,l,t,w,h,text,size=18,color=DARK,bold=False,align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,font=HEI,italic=False,lh=None):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    # 留白防贴边
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    for i,line in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        if lh: p.line_spacing=lh
        r=p.add_run(); r.text=line
        r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=color; r.font.name=font
    return tb

def rect(s,l,t,w,h,fill=None,line=None,line_w=1.0):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(line_w)
    sh.shadow.inherit=False; return sh

def rrect(s,l,t,w,h,fill=None,line=None,line_w=1.0):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(line_w)
    sh.shadow.inherit=False; return sh

def ln(s,x1,y1,x2,y2,color=NAVY,w=1.5):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,x1,y1,x2,y2)
    c.line.color.rgb=color; c.line.width=Pt(w); c.shadow.inherit=False; return c

def oval(s,l,t,w,h,fill=None,line=None,line_w=1.5):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,l,t,w,h)
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(line_w)
    sh.shadow.inherit=False; return sh

def pageno(s,n,total=16):
    txt(s,Inches(11.7),Inches(7.05),Inches(1.4),Inches(0.3),
        f"P. {n:02d}",size=10,color=GRAY,align=PP_ALIGN.RIGHT,font=SERIF)

def header(s,kicker,color=NAVY):
    rect(s,Inches(0.55),Inches(0.45),Inches(0.05),Inches(0.36),fill=color)
    txt(s,Inches(0.72),Inches(0.42),Inches(9),Inches(0.4),kicker,size=12,color=color,font=KAI)
    txt(s,Inches(9.6),Inches(0.42),Inches(3.1),Inches(0.4),
        "论文写作方法论 · 第一讲",size=10,color=GRAY,align=PP_ALIGN.RIGHT,font=SONG)

def hline(s,l,t,w,color=LINEC,w_pt=0.75):
    ln(s,l,t,l+w,t,color=color,w=w_pt)

# ========== 图表辅助 ==========
def style_axis(chart, color=GRAY, size=9):
    try:
        chart.category_axis.tick_labels.font.size = Pt(size)
        chart.category_axis.tick_labels.font.color.rgb = color
        chart.category_axis.format.line.color.rgb = LINEC
        if chart.has_value_axis:
            chart.value_axis.tick_labels.font.size = Pt(size)
            chart.value_axis.tick_labels.font.color.rgb = color
            chart.value_axis.format.line.color.rgb = LINEC
            chart.value_axis.major_gridlines.format.line.color.rgb = LINEC
    except Exception: pass

# ========== 简笔画 ==========
def draw_target(s,cx,cy,r,color=NAVY,accent=RUST):
    for ratio,lw in [(1.0,1.25),(0.72,1.25),(0.45,1.25)]:
        d=Emu(int(r*2*ratio))
        oval(s,cx-Emu(int(r*ratio)),cy-Emu(int(r*ratio)),d,d,line=color,line_w=lw)
    d=Emu(int(r*0.22))
    oval(s,cx-Emu(int(r*0.11)),cy-Emu(int(r*0.11)),d,d,fill=accent,line=accent,line_w=1.0)
    ln(s,cx-Emu(int(r*1.15)),cy,cx+Emu(int(r*1.15)),cy,color=color,w=0.75)
    ln(s,cx,cy-Emu(int(r*1.15)),cx,cy+Emu(int(r*1.15)),color=color,w=0.75)

def draw_magnifier(s,cx,cy,r,color=NAVY):
    oval(s,cx-Emu(int(r)),cy-Emu(int(r)),Emu(int(r*2)),Emu(int(r*2)),line=color,line_w=2.0)
    ln(s,cx+Emu(int(r*0.7)),cy+Emu(int(r*0.7)),cx+Emu(int(r*1.7)),cy+Emu(int(r*1.7)),color=color,w=3.0)

def draw_doc(s,l,t,w,h,color=NAVY,lines_n=4):
    rect(s,l,t,w,h,fill=PAPER,line=color,line_w=1.0)
    ln(s,l+w-Emu(int(w*0.2)),t,l+w,t+Emu(int(h*0.15)),color=color,w=1.0)
    ln(s,l+w-Emu(int(w*0.2)),t,l+w-Emu(int(w*0.2)),t+Emu(int(h*0.15)),color=color,w=0.75)
    ln(s,l+w-Emu(int(w*0.2)),t+Emu(int(h*0.15)),l+w,t+Emu(int(h*0.15)),color=color,w=0.75)
    for i in range(lines_n):
        ly=t+Emu(int(h*0.32))+Emu(int(h*0.15*i))
        end=0.78 if i%2==0 else 0.6
        ln(s,l+Emu(int(w*0.14)),ly,l+Emu(int(w*end)),ly,color=LINEC,w=0.75)

def draw_steps(s,l,t,bw,bh,n,color=NAVY,fill=CREAM):
    for i in range(n):
        h=Emu(int(bh*(i+1)/n)); x=l+Emu(int(bw*i))
        rect(s,x,t+bh-h,Emu(int(bw*0.92)),h,fill=fill,line=color,line_w=1.0)

def draw_person_sit(s,cx,cy,scale,color=NAVY):
    r=Emu(int(scale*0.5))
    oval(s,cx-Emu(int(r)),cy-Emu(int(r)),Emu(int(r*2)),Emu(int(r*2)),line=color,line_w=1.5)
    ln(s,cx,cy+r,cx-Emu(int(scale*0.5)),cy+Emu(int(scale*1.8)),color=color,w=1.5)
    ln(s,cx-Emu(int(scale*0.5)),cy+Emu(int(scale*1.8)),cx+Emu(int(scale*0.4)),cy+Emu(int(scale*1.9)),color=color,w=1.5)
    ln(s,cx,cy+Emu(int(scale*0.9)),cx+Emu(int(scale*1.2)),cy+Emu(int(scale*0.7)),color=color,w=1.5)

def draw_ruler(s,l,t,w,h,color=NAVY):
    rect(s,l,t,w,h,fill=CREAM,line=color,line_w=1.0)
    n=8
    for i in range(n+1):
        x=l+Emu(int(w*i/n)); tl=h*0.4 if i%2==0 else h*0.25
        ln(s,x,t,x,t+Emu(int(tl)),color=color,w=0.75)

# ===================================================================
# P01 封面
# ===================================================================
s=slide()
# 左侧：羽毛笔 + 文档堆
px,py=Inches(1.1),Inches(2.0)
ln(s,px,py,px+Inches(1.7),py+Inches(1.7),color=NAVY,w=2.0)
ln(s,px,py,px-Inches(0.18),py+Inches(0.28),color=NAVY,w=2.5)
ln(s,px-Inches(0.18),py+Inches(0.28),px+Inches(0.05),py+Inches(0.1),color=NAVY,w=1.0)
for i in range(3):
    ln(s,px+Inches(0.4*i),py+Inches(0.4*i),px+Inches(0.4*i)+Inches(0.25),py+Inches(0.4*i)-Inches(0.22),color=GRAY,w=0.75)
for i in range(3):
    draw_doc(s,Inches(1.0)+Inches(0.06*i),Inches(4.4)+Inches(0.08*i),Inches(1.5),Inches(1.0),color=GRAY,lines_n=3)

txt(s,Inches(4.2),Inches(0.9),Inches(8),Inches(0.5),"论文写作方法论 · 第一讲",size=15,color=RUST,font=KAI)
hline(s,Inches(4.2),Inches(1.45),Inches(3.0),color=RUST,w_pt=1.25)
txt(s,Inches(4.2),Inches(2.0),Inches(8.5),Inches(2.2),
    "毕业论文到底\n在考你什么？",size=46,color=INK,bold=True,font=SONG)
txt(s,Inches(4.2),Inches(4.35),Inches(8.5),Inches(0.8),
    "90% 的人，从第一步就错了",size=24,color=RUST,font=KAI)
hline(s,Inches(0.6),Inches(6.6),Inches(12.1),color=LINEC,w_pt=1.0)
txt(s,Inches(0.6),Inches(6.75),Inches(8),Inches(0.4),
    "写论文之前，先厘清这件事",size=13,color=GRAY,font=KAI)
txt(s,Inches(9.5),Inches(6.75),Inches(3.2),Inches(0.4),
    "Lecture 01 / 100",size=12,color=GRAY,align=PP_ALIGN.RIGHT,font=SERIF)
pageno(s,1)

# ===================================================================
# P02 四类困境 + 现状数据柱图
# ===================================================================
s=slide()
header(s,"引　言　　写作之困")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "动笔之前，先确认下列困境是否常见：",size=23,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 左侧：四条困境（紧凑分区 1.9-6.4）
plights=[
    ("1","文档新建、标题已定，却面对光标久坐，难以下笔首句。"),
    ("2","积稿三万字，导师以\"选题存疑\"为由要求推倒重来。"),
    ("3","查重率超出阈值，降重耗时反多于撰写本身。"),
    ("4","答辩席上被问及\"所研究为何\"，竟一时语塞。"),
]
y=Inches(1.95)
for num,c in plights:
    rect(s,Inches(0.85),y+Inches(0.05),Inches(0.04),Inches(0.7),fill=RUST)
    txt(s,Inches(1.0),y+Inches(0.02),Inches(0.45),Inches(0.7),num,size=15,color=RUST,bold=True,anchor=MSO_ANCHOR.MIDDLE,font=SERIF)
    txt(s,Inches(1.55),y,Inches(5.7),Inches(0.75),c,size=14,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE,lh=1.1)
    y+=Inches(1.1)
txt(s,Inches(0.72),Inches(6.5),Inches(6.4),Inches(0.4),
    "上述困境多源于初始认知偏差，而非写作功力不足。",
    size=13,color=RUST,font=KAI)

# 右侧：柱状图——论文写作常见困难分布（示意性数据）
txt(s,Inches(7.5),Inches(1.95),Inches(5.2),Inches(0.4),
    "图1　学生论文写作常见困难分布",size=12,color=INK,bold=True,font=SONG)
txt(s,Inches(7.5),Inches(2.3),Inches(5.2),Inches(0.3),
    "（示意性数据，基于教学观察）",size=10,color=GRAY,font=KAI)
cd=CategoryChartData()
cd.categories=["选题困惑","结构混乱","查重超标","格式错误","文献不足","答辩紧张"]
cd.add_series("占比",(38,27,22,19,18,15))
gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(7.5),Inches(2.7),Inches(5.3),Inches(3.6),cd).chart
gf.has_title=False; gf.has_legend=False
plot=gf.plots[0]; plot.gap_width=80
ser=plot.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb=NAVY
style_axis(gf)
pageno(s,2)

# ===================================================================
# P03 核心论点 + 误区对比
# ===================================================================
s=slide()
header(s,"本　论　　论文的本质")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "一个常被忽视的根本区分：",size=23,color=INK,bold=True,font=SONG)

# 左：误区（0.8-6.5）
rect(s,Inches(0.8),Inches(1.95),Inches(5.7),Inches(0.55),fill=RUST)
txt(s,Inches(0.8),Inches(1.95),Inches(5.7),Inches(0.55),"认识误区",
    size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
rrect(s,Inches(0.8),Inches(2.5),Inches(5.7),Inches(2.3),fill=CREAM,line=LINEC,line_w=0.75)
txt(s,Inches(1.05),Inches(2.65),Inches(5.3),Inches(2.0),
    "论文 = 一篇长文章\n\n认为字数凑足、排版工整、\n查重过关，即算完成。",
    size=16,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE,lh=1.15)
draw_doc(s,Inches(5.4),Inches(2.8),Inches(0.8),Inches(1.05),color=RUST,lines_n=4)

# 右：正解（6.83-11.8）
rect(s,Inches(6.83),Inches(1.95),Inches(5.7),Inches(0.55),fill=NAVY)
txt(s,Inches(6.83),Inches(1.95),Inches(5.7),Inches(0.55),"本质正解",
    size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
rrect(s,Inches(6.83),Inches(2.5),Inches(5.7),Inches(2.3),fill=PAPER,line=NAVY,line_w=1.0)
txt(s,Inches(7.08),Inches(2.65),Inches(5.3),Inches(2.0),
    "论文 = 一次完整的\n学术研究训练\n\n所考者为\"做研究\"，\n而非\"写文章\"。",
    size=16,color=INK,bold=True,font=KAI,anchor=MSO_ANCHOR.MIDDLE,lh=1.15)
draw_target(s,Inches(11.6),Inches(3.4),Inches(0.5),color=NAVY,accent=RUST)

# 底部金句（精确区间 5.0-6.9）
rect(s,Inches(0.8),Inches(5.05),Inches(11.73),Inches(1.6),fill=INK)
txt(s,Inches(1.1),Inches(5.15),Inches(11.2),Inches(0.4),
    "核心命题",size=12,color=RGBColor(0xC8,0xC0,0xB0),font=KAI)
txt(s,Inches(1.1),Inches(5.55),Inches(11.2),Inches(1.0),
    "论文并非\"写\"出来的，而是\"研究\"出来的 —— 写作只是研究终局的呈现环节。",
    size=19,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER,font=SONG)
pageno(s,3)

# ===================================================================
# P04 三种核心能力 + 能力权重饼图
# ===================================================================
s=slide()
header(s,"本　论　　所考的三种能力")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "论文所考核的三种核心能力，缺一不可：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 左侧三段（0.85-7.0）
items=[
    ("01","发现问题","能在一领域内，找到一个值得研究且可研究的问题。",NAVY),
    ("02","解决问题","能以相宜的方法，切实验证并回答这一问题。",RUST),
    ("03","表达问题","能依学术规范，将研究过程与结论清晰呈现于人。",OLIVE),
]
y=Inches(2.05)
for num,t,d,col in items:
    rect(s,Inches(0.85),y,Inches(0.05),Inches(1.35),fill=col)
    txt(s,Inches(1.05),y-Inches(0.02),Inches(0.7),Inches(0.5),num,size=16,color=col,bold=True,font=SERIF)
    txt(s,Inches(1.8),y-Inches(0.02),Inches(3.0),Inches(0.5),t,size=17,color=INK,bold=True,font=SONG)
    txt(s,Inches(1.8),y+Inches(0.45),Inches(5.0),Inches(0.9),d,size=13,color=DARK,font=KAI,lh=1.1)
    y+=Inches(1.45)
txt(s,Inches(0.85),Inches(6.55),Inches(6.3),Inches(0.4),
    "注：写作能力仅属第三项之子集，而多数人将全部精力倾注于此。",
    size=12,color=RUST,font=KAI)

# 右侧饼图：能力权重分配（示意）
txt(s,Inches(7.5),Inches(2.0),Inches(5.2),Inches(0.4),
    "图2　三种能力的考核权重（示意）",size=12,color=INK,bold=True,font=SONG)
cd=CategoryChartData()
cd.categories=["发现问题","解决问题","表达问题"]
cd.add_series("权重",(40,35,25))
pf=s.shapes.add_chart(XL_CHART_TYPE.PIE,Inches(7.4),Inches(2.5),Inches(5.4),Inches(3.6),cd).chart
pf.has_title=False
pf.has_legend=True; pf.legend.position=XL_LEGEND_POSITION.RIGHT; pf.legend.include_in_layout=False
pf.legend.font.size=Pt(11); pf.legend.font.color.rgb=DARK
plot=pf.plots[0]; plot.has_data_labels=True
dl=plot.data_labels; dl.number_format='0%'; dl.number_format_is_linked=False
dl.font.size=Pt(11); dl.font.color.rgb=WHITE; dl.font.bold=True
# 上色
cols=[NAVY,RUST,OLIVE]
for i,pt in enumerate(plot.series[0].points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb=cols[i]
pageno(s,4)

# ===================================================================
# P05 第一步 = 问题意识 + 纠错成本柱图
# ===================================================================
s=slide()
header(s,"本　论　　何谓\"第一步\"")
txt(s,Inches(0.72),Inches(1.05),Inches(7),Inches(0.55),
    "所谓\"第一步\"，并非打开 Word —— 而是：",size=20,color=INK,bold=True,font=SONG)
rect(s,Inches(0.8),Inches(1.75),Inches(6.9),Inches(1.25),fill=INK)
txt(s,Inches(1.05),Inches(1.75),Inches(6.4),Inches(1.25),
    "建立「问题意识」\n想清楚 —— 此文欲回答之问题，究竟为何？",
    size=18,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE,font=SONG,lh=1.15)

# 流程条（0.8-7.7, y=3.35）
txt(s,Inches(0.8),Inches(3.25),Inches(7),Inches(0.3),
    "若径跳此步、埋头即写，则后续各阶段均\"带病推进\"：",size=12,color=GRAY,font=KAI)
fl=["选题","读文献","开题","写正文","答辩"]; fc=[NAVY,NAVY,NAVY,GRAY,GRAY]
fx=Inches(0.8)
for i in range(5):
    rrect(s,fx,Inches(3.6),Inches(1.2),Inches(0.55),fill=PAPER,line=fc[i],line_w=1.25)
    txt(s,fx,Inches(3.6),Inches(1.2),Inches(0.55),fl[i],size=13,color=fc[i],bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
    if i<4:
        txt(s,fx+Inches(1.2),Inches(3.6),Inches(0.35),Inches(0.55),"→",size=14,color=GRAY,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SERIF)
    fx+=Inches(1.55)
txt(s,Inches(0.8),Inches(4.35),Inches(6.9),Inches(0.9),
    "问题未明，则愈往后纠错代价愈昂 ——\n此即\"90% 自第一步误\"之真义。",
    size=15,color=DARK,bold=True,font=KAI,lh=1.2)

# 右侧：纠错成本随阶段上升柱图（强说服力数据图）
txt(s,Inches(8.0),Inches(1.05),Inches(4.9),Inches(0.4),
    "图3　不同阶段发现错误的纠错成本",size=12,color=INK,bold=True,font=SONG)
txt(s,Inches(8.0),Inches(1.4),Inches(4.9),Inches(0.3),
    "（相对成本指数，示意）",size=10,color=GRAY,font=KAI)
cd=CategoryChartData()
cd.categories=["选题期","开题期","写作期","答辩期"]
cd.add_series("成本",(1,3,8,18))
cf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(8.0),Inches(1.8),Inches(4.8),Inches(3.4),cd).chart
cf.has_title=False; cf.has_legend=False
plot=cf.plots[0]; plot.gap_width=60
ser=plot.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb=RUST
style_axis(cf)
txt(s,Inches(8.0),Inches(5.25),Inches(4.9),Inches(0.7),
    "结论：错误发现越晚，\n推倒重来的代价呈倍数级增长。",
    size=13,color=RUST,bold=True,font=KAI,lh=1.15)
pageno(s,5)

# ===================================================================
# P06-08 三个案例（含简笔画，分区明确）
# ===================================================================
def case_slide(n,tag,title,scene,result,root,color,draw_fn):
    s=slide()
    header(s,f"案　例　　{tag}",color)
    txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),title,size=25,color=color,bold=True,font=SONG)
    hline(s,Inches(0.72),Inches(1.7),Inches(11.9))
    # 左：简笔区 0.8-5.0
    txt(s,Inches(0.8),Inches(1.95),Inches(4.0),Inches(0.3),"〔情形示意〕",size=11,color=GRAY,font=KAI,align=PP_ALIGN.CENTER)
    draw_fn(s)
    # 右：文字区 5.2-12.5
    txt(s,Inches(5.2),Inches(1.95),Inches(3),Inches(0.35),"〔场景〕",size=13,color=color,bold=True,font=SONG)
    txt(s,Inches(5.2),Inches(2.35),Inches(7.4),Inches(1.5),scene,size=14,color=DARK,font=KAI,lh=1.15)
    txt(s,Inches(5.2),Inches(3.9),Inches(3),Inches(0.35),"〔结果〕",size=13,color=RUST,bold=True,font=SONG)
    txt(s,Inches(5.2),Inches(4.3),Inches(7.4),Inches(1.3),result,size=14,color=DARK,font=KAI,lh=1.15)
    # 底部病根 5.9-6.95
    rect(s,Inches(0.8),Inches(5.95),Inches(11.73),Inches(0.95),fill=color)
    txt(s,Inches(1.1),Inches(5.95),Inches(1.8),Inches(0.95),"病根",size=14,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
    ln(s,Inches(2.6),Inches(6.1),Inches(2.6),Inches(6.75),color=WHITE,w=0.75)
    txt(s,Inches(2.85),Inches(5.95),Inches(9.4),Inches(0.95),root,size=15,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE,font=KAI,lh=1.1)
    pageno(s,n)

def draw_caseA(s):
    draw_doc(s,Inches(1.3),Inches(2.4),Inches(1.6),Inches(2.2),color=NAVY,lines_n=0)
    txt(s,Inches(1.3),Inches(3.15),Inches(1.6),Inches(0.4),"（空白文档）",size=11,color=GRAY,align=PP_ALIGN.CENTER,font=KAI)
    draw_person_sit(s,Inches(3.7),Inches(2.6),Inches(1.1),color=NAVY)
    oval(s,Inches(3.3),Inches(2.1),Inches(0.18),Inches(0.18),fill=RUST,line=RUST)
case_slide(6,"甲","未及深思，仓促动笔",
    "张某以\"论新媒体之发展\"为题，经导师一句\"先写\"，便埋头撰成三万字。",
    "中期答辩被问：所研究之具体问题为何？竟不能答。\n终推倒重写，延宕半年。",
    "题目是\"话题\"而非\"问题\" —— 自始缺乏待答之疑问，写作遂失靶心。",
    RUST,draw_caseA)

def draw_caseB(s):
    oval(s,Inches(1.1),Inches(2.4),Inches(3.2),Inches(2.6),line=RUST,line_w=1.5)
    txt(s,Inches(1.1),Inches(2.45),Inches(3.2),Inches(0.35),"选题范围过大",size=11,color=RUST,align=PP_ALIGN.CENTER,font=KAI)
    pts=[(1.9,3.2),(2.7,3.1),(3.3,3.6),(2.1,4.0),(3.1,4.2),(2.5,4.5),(1.7,3.7),(3.4,4.0)]
    for px,py in pts:
        oval(s,Inches(px),Inches(py),Inches(0.15),Inches(0.15),fill=NAVY,line=NAVY)
case_slide(7,"乙","选题过大，力有不逮",
    "李某欲撰\"中国教育公平问题研究\"，资料愈积愈多，每一点皆足成一书。",
    "欲写者众，能深者寡。评委评曰：\n此文近于综述，未见独立研究。",
    "选题颗粒度过大 —— 以\"整个领域\"充当\"一个问题\"，超出个人研究承载力。",
    RGBColor(0xA8,0x6A,0x3A),draw_caseB)

def draw_caseC(s):
    oval(s,Inches(1.3),Inches(3.2),Inches(1.5),Inches(1.5),line=NAVY,line_w=1.5)
    oval(s,Inches(2.3),Inches(3.2),Inches(1.5),Inches(1.5),line=RUST,line_w=1.5)
    txt(s,Inches(0.9),Inches(2.4),Inches(2.0),Inches(0.35),"本例选题",size=11,color=NAVY,align=PP_ALIGN.CENTER,font=KAI)
    txt(s,Inches(2.6),Inches(4.8),Inches(2.0),Inches(0.35),"前人已究",size=11,color=RUST,align=PP_ALIGN.CENTER,font=KAI)
case_slide(8,"丙","未察文献，选题撞车",
    "王某自以为得一绝佳视角，径自开干，撰毕方知 —— 三年前已有人做同一研究。",
    "创新性归零，等于白做。\n所有结论皆被\"前人已证\"所否。",
    "缺乏文献意识 —— 不知他人所做为何，\"自以为新\"实乃重复劳动。",
    NAVY,draw_caseC)

# ===================================================================
# P09 共同死因 + 失败阶段分布饼图
# ===================================================================
s=slide()
header(s,"归　纳　　三案同源")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "三案异途，然同归于一根本之因：",size=23,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))

# 左：汇聚简笔 0.8-4.3
cx,cy=Inches(2.3),Inches(4.4)
ln(s,Inches(1.0),Inches(2.5),cx,cy,color=RUST,w=1.25)
ln(s,Inches(1.0),Inches(4.4),cx,cy,color=RGBColor(0xA8,0x6A,0x3A),w=1.25)
ln(s,Inches(1.0),Inches(6.0),cx,cy,color=NAVY,w=1.25)
oval(s,cx-Inches(0.14),cy-Inches(0.14),Inches(0.28),Inches(0.28),fill=INK,line=INK)
txt(s,Inches(1.0),Inches(2.25),Inches(1.2),Inches(0.3),"案甲",size=12,color=RUST,font=KAI)
txt(s,Inches(1.0),Inches(4.15),Inches(1.2),Inches(0.3),"案乙",size=12,color=RGBColor(0xA8,0x6A,0x3A),font=KAI)
txt(s,Inches(1.0),Inches(5.75),Inches(1.2),Inches(0.3),"案丙",size=12,color=NAVY,font=KAI)
txt(s,Inches(2.6),Inches(4.25),Inches(1.7),Inches(0.3),"同一病根",size=12,color=INK,bold=True,font=KAI)

# 右上：归因表 4.5-12.5
rows=[("案甲","题目是话题，非问题","缺 · 问题意识"),
      ("案乙","选题过大，力有不逮","缺 · 问题意识"),
      ("案丙","创新已被前人验证","缺 · 文献意识")]
y=Inches(2.1)
for lt,mt,rt in rows:
    rrect(s,Inches(4.5),y,Inches(4.4),Inches(0.95),fill=CREAM,line=LINEC,line_w=0.75)
    txt(s,Inches(4.7),y,Inches(0.9),Inches(0.95),lt,size=14,color=INK,bold=True,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
    ln(s,Inches(5.6),y+Inches(0.15),Inches(5.6),y+Inches(0.8),color=LINEC,w=0.75)
    txt(s,Inches(5.8),y,Inches(3.0),Inches(0.95),mt,size=13,color=DARK,anchor=MSO_ANCHOR.MIDDLE,font=KAI)
    rect(s,Inches(9.0),y,Inches(3.53),Inches(0.95),fill=INK)
    txt(s,Inches(9.0),y,Inches(3.53),Inches(0.95),rt,size=14,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
    y+=Inches(1.1)

# 右下：论文失败环节分布饼图
txt(s,Inches(4.5),Inches(5.5),Inches(8),Inches(0.35),
    "图4　论文失败的关键环节分布（示意）",size=12,color=INK,bold=True,font=SONG)
cd=CategoryChartData()
cd.categories=["选题阶段","开题阶段","写作阶段","其他"]
cd.add_series("占比",(62,16,15,7))
pf=s.shapes.add_chart(XL_CHART_TYPE.PIE,Inches(7.6),Inches(5.5),Inches(5.2),Inches(1.4),cd).chart
pf.has_title=False
pf.has_legend=True; pf.legend.position=XL_LEGEND_POSITION.RIGHT; pf.legend.include_in_layout=False
pf.legend.font.size=Pt(9); pf.legend.font.color.rgb=DARK
plot=pf.plots[0]; plot.has_data_labels=True
dl=plot.data_labels; dl.number_format='0%'; dl.number_format_is_linked=False
dl.font.size=Pt(9); dl.font.color.rgb=WHITE; dl.font.bold=True
cols=[RUST,NAVY,RGBColor(0xA8,0x6A,0x3A),GRAY]
for i,pt in enumerate(plot.series[0].points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb=cols[i]
txt(s,Inches(0.72),Inches(6.7),Inches(12),Inches(0.4),
    "结论：约六成论文败局，败于动笔之前 —— 即败于\"未明所研为何\"。",
    size=14,color=RUST,bold=True,font=SONG)
pageno(s,9)

# ===================================================================
# P10-12 三大认知
# ===================================================================
def cognition_slide(n,idx,title,wrong,right,key,color,draw_fn):
    s=slide()
    header(s,f"认　知　　其{idx}",color)
    txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),title,size=24,color=color,bold=True,font=SONG)
    hline(s,Inches(0.72),Inches(1.7),Inches(11.9))
    # 左简笔 0.8-4.6
    txt(s,Inches(0.8),Inches(1.95),Inches(3.7),Inches(0.3),"〔示意〕",size=11,color=GRAY,font=KAI,align=PP_ALIGN.CENTER)
    draw_fn(s)
    # 右：误区/正解 5.0-12.5
    txt(s,Inches(5.0),Inches(2.1),Inches(3.5),Inches(0.35),"〔误区〕",size=13,color=RUST,bold=True,font=SONG)
    rrect(s,Inches(5.0),Inches(2.5),Inches(3.4),Inches(1.7),fill=CREAM,line=LINEC,line_w=0.75)
    txt(s,Inches(5.2),Inches(2.5),Inches(3.05),Inches(1.7),wrong,size=13,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE,lh=1.15)
    txt(s,Inches(8.9),Inches(2.1),Inches(3.6),Inches(0.35),"〔正解〕",size=13,color=OLIVE,bold=True,font=SONG)
    rrect(s,Inches(8.9),Inches(2.5),Inches(3.6),Inches(1.7),fill=PAPER,line=OLIVE,line_w=1.0)
    txt(s,Inches(9.1),Inches(2.5),Inches(3.2),Inches(1.7),right,size=13,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE,lh=1.15)
    # 底部金句 5.2-6.9
    rect(s,Inches(0.8),Inches(5.15),Inches(11.73),Inches(1.55),fill=color)
    txt(s,Inches(1.1),Inches(5.25),Inches(11.2),Inches(0.35),"〔核心要义〕",size=11,color=RGBColor(0xD8,0xD0,0xC0),font=KAI)
    txt(s,Inches(1.1),Inches(5.6),Inches(11.2),Inches(1.0),key,size=18,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER,font=SONG)
    pageno(s,n)

def draw_c1(s):
    draw_target(s,Inches(2.6),Inches(3.6),Inches(1.4),color=NAVY,accent=RUST)
cognition_slide(10,"一","认知一 · 问题意识 —— 论文须有\"靶心\"",
    "以\"话题\"充当\"问题\"。\n\"论人工智能\"是话题，\n并无待答之疑问。",
    "将话题收束为一个疑问。\n\"AI客服能否提升用户信任？\"\n此为可验证之问题。",
    "好论文回答一个具体的问题，\n而非泛论一个宽泛的话题。",
    NAVY,draw_c1)

def draw_c2(s):
    draw_magnifier(s,Inches(2.5),Inches(3.2),Inches(1.15),color=RGBColor(0xA8,0x6A,0x3A))
    for i,x in enumerate([1.2,2.3,3.4]):
        draw_doc(s,Inches(x),Inches(4.6),Inches(0.9),Inches(0.7),color=GRAY,lines_n=3)
cognition_slide(11,"二","认知二 · 文献意识 —— 立于前人之肩",
    "闭门造车，凭直觉选题，\n自以为有所创新。",
    "先查清他人已做到何处，\n找到\"未做\"或\"未足\"之缝隙，\n那才是自己的位置。",
    "创新不是凭空而来，而是读出来的 ——\n文献之中，藏有选题的答案。",
    RGBColor(0xA8,0x6A,0x3A),draw_c2)

def draw_c3(s):
    draw_ruler(s,Inches(1.2),Inches(3.3),Inches(3.0),Inches(0.7),color=OLIVE)
    txt(s,Inches(1.2),Inches(4.25),Inches(3.0),Inches(0.35),"学术规范",size=12,color=OLIVE,bold=True,align=PP_ALIGN.CENTER,font=SONG)
    draw_doc(s,Inches(1.9),Inches(4.8),Inches(1.6),Inches(1.1),color=NAVY,lines_n=4)
cognition_slide(12,"三","认知三 · 规范意识 —— 循规矩以行事",
    "以为内容好即可，\n格式、引用、查重皆为细枝末节，\n末了草草了事。",
    "学术规范是\"入场券\"，\n不合规即一票否决，\n内容虽佳亦属徒然。",
    "规范不是束缚，而是通行证 ——\n先求合规，再论优劣。",
    OLIVE,draw_c3)

# ===================================================================
# P13 良窳之辨（学术三线表）
# ===================================================================
s=slide()
header(s,"判　准　　良窳之辨")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "判一篇论文之优劣，可观此四端：",size=23,color=INK,bold=True,font=SONG)
tl=Inches(0.8); tw=Inches(11.73); ty=Inches(2.15); rh=Inches(0.82)
hline(s,tl,ty,tw,color=INK,w_pt=1.5)
heads=["维　度","窳者（差）","良者（好）"]
hx=[tl,tl+Inches(2.6),tl+Inches(6.6)]; hw=[Inches(2.6),Inches(4.0),Inches(5.13)]
for i,h in enumerate(heads):
    txt(s,hx[i],ty,hw[i],rh,h,size=16,color=INK,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
hline(s,tl,ty+rh,tw,color=INK,w_pt=0.75)
rows=[("问　题","话题宽泛，无待答之疑问","聚焦一个具体可答的问题"),
      ("文　献","罗列堆砌，近似流水账","有述有评，引出自身研究"),
      ("方　法","方法与问题相脱节","方法服务于回答问题"),
      ("结　论","复述常识，无知识增量","有自身发现，虽小亦珍")]
ry=ty+rh
for a,b,c in rows:
    txt(s,hx[0],ry,hw[0],rh,a,size=15,color=DARK,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SONG)
    txt(s,hx[1]+Inches(0.2),ry,hw[1]-Inches(0.2),rh,b,size=14,color=DARK,anchor=MSO_ANCHOR.MIDDLE,font=KAI)
    txt(s,hx[2]+Inches(0.2),ry,hw[2]-Inches(0.2),rh,c,size=14,color=INK,anchor=MSO_ANCHOR.MIDDLE,font=KAI)
    ry+=rh
hline(s,tl,ry,tw,color=INK,w_pt=1.5)
txt(s,Inches(0.72),Inches(6.85),Inches(12),Inches(0.4),
    "四端皆过，即为合格之文；良者，不过各端更为扎实而已。",
    size=13,color=GRAY,font=KAI)
pageno(s,13)

# ===================================================================
# P14 自查六问
# ===================================================================
s=slide()
header(s,"自　测　　动笔前之六问",RUST)
txt(s,Inches(0.72),Inches(1.0),Inches(12),Inches(0.55),
    "动笔之前，请以下列六问自省：",size=23,color=INK,bold=True,font=SONG)
txt(s,Inches(0.72),Inches(1.58),Inches(12),Inches(0.3),
    "〔建议截图保存〕",size=11,color=RUST,font=KAI)
hline(s,Inches(0.72),Inches(2.0),Inches(11.9))
checks=[
    "研究问题是一个\"疑问\"还是\"话题\"？能否用一句话讲清？",
    "此问题，前人是否研究过？已做到何种程度？",
    "以当前的时间、能力、资源，能否完成此项研究？",
    "将采用何种方法回答此问题？方法与问题是否匹配？",
    "研究意义为何？理论意义与实践意义分别是什么？",
    "是否已了解本校的格式、字数、查重、送审等要求？",
]
y=Inches(2.35)
for i,c in enumerate(checks):
    oval(s,Inches(0.9),y+Inches(0.08),Inches(0.42),Inches(0.42),fill=PAPER,line=INK,line_w=1.0)
    txt(s,Inches(0.9),y+Inches(0.08),Inches(0.42),Inches(0.42),str(i+1),size=14,color=INK,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,font=SERIF)
    rect(s,Inches(1.5),y+Inches(0.08),Inches(0.04),Inches(0.42),fill=RUST)
    txt(s,Inches(1.7),y,Inches(10.8),Inches(0.58),c,size=15,color=DARK,font=KAI,anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(0.68)
txt(s,Inches(0.72),Inches(6.7),Inches(12),Inches(0.4),
    "六问皆能答，方可谓之\"已想清楚\" —— 而后方可动笔。",
    size=13,color=RUST,bold=True,font=SONG)
pageno(s,14)

# ===================================================================
# P15 全流程甘特图
# ===================================================================
s=slide()
header(s,"展　望　　全流程时序")
txt(s,Inches(0.72),Inches(1.05),Inches(12),Inches(0.55),
    "论文非一蹴而就，乃一岁之长征：",size=22,color=INK,bold=True,font=SONG)
hline(s,Inches(0.72),Inches(1.75),Inches(11.9))
months=["1月","2月","3月","4月","5月","6月","7月","8月","9月","10月","11月","12月"]
ml=Inches(2.4); mw=Inches(10.0)
for i,m in enumerate(months):
    x=ml+Emu(int(mw*i/11))
    txt(s,x-Inches(0.25),Inches(2.1),Inches(0.5),Inches(0.3),m,size=10,color=GRAY,align=PP_ALIGN.CENTER,font=SERIF)
    ln(s,x,Inches(2.45),x,Inches(2.55),color=LINEC,w=0.75)
hline(s,ml,Inches(2.5),mw,color=INK,w_pt=1.0)
phases=[("选题与文献",0,2,NAVY),("开题报告",2,1,NAVY),
        ("研究实施",3,5,RGBColor(0xA8,0x6A,0x3A)),("正文写作",6,4,RGBColor(0xA8,0x6A,0x3A)),
        ("修改查重",9,2,OLIVE),("答辩定稿",11,1,OLIVE)]
yy=Inches(2.8)
for name,start,span,col in phases:
    txt(s,Inches(0.6),yy,Inches(1.7),Inches(0.4),name,size=13,color=DARK,bold=True,align=PP_ALIGN.RIGHT,font=SONG,anchor=MSO_ANCHOR.MIDDLE)
    bx=ml+Emu(int(mw*start/11)); bw=Emu(int(mw*span/11))
    rect(s,bx,yy-Inches(0.05),bw,Inches(0.35),fill=col)
    yy+=Inches(0.55)
txt(s,Inches(0.72),Inches(6.2),Inches(12),Inches(0.8),
    "要旨：前段多费时，则后段少走弯路。\n后续各阶段，皆将另设专讲详述。",
    size=15,color=DARK,bold=True,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER,font=KAI)
pageno(s,15)

# ===================================================================
# P16 结尾 + 下期预告
# ===================================================================
s=slide()
header(s,"结　语")
txt(s,Inches(0.72),Inches(1.2),Inches(12),Inches(0.5),"本讲核心命题",size=15,color=RUST,font=KAI)
hline(s,Inches(0.72),Inches(1.8),Inches(3.0),color=RUST,w_pt=1.25)
txt(s,Inches(0.72),Inches(2.1),Inches(12),Inches(1.2),
    "论文并非\"写\"出来的，而是\"研究\"出来的 ——\n先想清楚欲答何问，而后动笔。",
    size=24,color=INK,bold=True,font=SONG,lh=1.2)
txt(s,Inches(0.72),Inches(4.05),Inches(12),Inches(0.4),"下期预告",size=14,color=RUST,font=KAI)
hline(s,Inches(0.72),Inches(4.5),Inches(11.9))
txt(s,Inches(0.72),Inches(4.75),Inches(12),Inches(0.4),"第　二　讲",size=13,color=GRAY,font=SONG)
txt(s,Inches(0.72),Inches(5.15),Inches(12),Inches(1.0),
    "本科 / 硕士 / 博士论文的本质区别\n—— 勿以本科之思维，撰硕士之论文",
    size=21,color=INK,bold=True,font=SONG)
hline(s,Inches(0.6),Inches(6.85),Inches(12.1))
txt(s,Inches(0.6),Inches(6.95),Inches(12),Inches(0.4),
    "若觉有益，望垂注之 —— 期期相续，将论文一事，讲之通透。",
    size=12,color=GRAY,font=KAI)
pageno(s,16)

out=r"D:\project\code2\temp_pastemd\lw\1期\毕业论文到底在考你什么.ppt"
prs.save(out)
print("已保存:",out)
print("总页数:",len(prs.slides._sldIdLst))